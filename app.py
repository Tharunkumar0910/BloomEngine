import sys
if __name__ == "__main__":
    sys.modules['app'] = sys.modules['__main__']

import os
import uuid
import csv
import io
import json
import time
import re
from datetime import datetime
import threading
import functools
import torch
import pandas as pd
import platform
import transformers
from flask import Flask, request, jsonify, render_template, session, send_file, Response
from transformers import (
    AutoModelForSequenceClassification,
    AutoTokenizer,
    AutoModelForSeq2SeqLM,
)
import PyPDF2
import docx
from fpdf import FPDF
from werkzeug.utils import secure_filename
import pptx
from dataclasses import dataclass
from config import (
    MODES,
    ACTIVE_MODE,
    DOMAIN_MAP,
    NUM_CANDIDATES_BY_BLOOM,
)
from prompt_templates import build_prompt
from validation_engine import evaluate_candidate
from candidate_ranker import calculate_nlp_rank_score, rank_candidates, rank_candidates_dicts
from bloom_validator import validate_bloom_verbs
from spacy_utils import normalize_embedding_key




app = Flask(__name__)
app.secret_key = "super_secret_bloom_key_123"
app.config["UPLOAD_FOLDER"] = "static/uploads"

# Ensure upload directory exists
os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)

# In-memory session storage for large batch results
MEMORY_STORE = {}
BATCH_HISTORY = []
SESSION_STATE = {
    "manual_classifications": [],
    "batch_history": BATCH_HISTORY,
}


# Startup timestamp
STARTUP_TIMESTAMP = time.time()

# Performance Logging Flags
ENABLE_PERFORMANCE_LOGS = True
ENABLE_GENERATION_PERFORMANCE_LOGS = True

# Thread locks
MODEL_LOCK = threading.Lock()
EMBEDDING_CACHE_LOCK = threading.Lock()

# Embedding Cache & Stats
EMBEDDING_CACHE = {}
EMBEDDING_CACHE_HITS = 0
EMBEDDING_CACHE_MISSES = 0

# Warmup state -- set to True only after models load AND a test inference passes
_warmup_succeeded = False
_load_errors: list = []


# Helper to get memory usage of current process
def get_current_process_memory():
    try:
        import os
        import subprocess
        pid = os.getpid()
        cmd = f"powershell -Command \"(Get-Process -Id {pid}).WorkingSet64\""
        out = subprocess.check_output(cmd, shell=True).decode().strip()
        return float(out) / (1024 * 1024)  # MB
    except Exception:
        return 0.0


# Helper to get cached embedding
def get_cached_embedding(text: str, st_model):
    global EMBEDDING_CACHE_HITS, EMBEDDING_CACHE_MISSES
    key = normalize_embedding_key(text)
    with EMBEDDING_CACHE_LOCK:
        if key in EMBEDDING_CACHE:
            EMBEDDING_CACHE_HITS += 1
            return EMBEDDING_CACHE[key]

    # Compute embedding outside of lock
    emb = st_model.encode(key, convert_to_tensor=True)

    with EMBEDDING_CACHE_LOCK:
        EMBEDDING_CACHE[key] = emb
        EMBEDDING_CACHE_MISSES += 1
        return emb



# Global model references
deberta_model = None
deberta_tokenizer = None
flan_model = None
flan_tokenizer = None
_st_model = None


def _get_st_model():
    """Return the globally cached sentence-transformer model."""
    global _st_model
    if _st_model is None:
        with MODEL_LOCK:
            if _st_model is None:
                try:
                    from sentence_transformers import SentenceTransformer
                    print("[LAZY] Loading SentenceTransformer (all-MiniLM-L6-v2) on demand...")
                    _st_model = SentenceTransformer("all-MiniLM-L6-v2")
                except Exception as exc:  # pragma: no cover
                    print(f"[WARNING] sentence-transformers unavailable: {exc}")
                    _st_model = None
    return _st_model


DEBERTA_PATH = "./deberta_bloom_model"
FLAN_PATH = "./flan_t5_model"

# Mappings — use Easy / Medium / Hard to match FLAN-T5 fine-tuning labels.
# The model was trained with these exact strings; using Moderate/Difficult
# causes an inference/training mismatch that prevents accepted generations.
DIFFICULTY_TO_BLOOM = {
    "Easy": ["Remember", "Understand"],
    "Medium": ["Apply", "Analyze"],
    "Hard": ["Evaluate", "Create"],
}

BLOOM_TO_DIFFICULTY = {
    "Remember": "Easy",
    "Understand": "Easy",
    "Apply": "Medium",
    "Analyze": "Medium",
    "Evaluate": "Hard",
    "Create": "Hard",
}


@dataclass
class ValidationResult:
    generated_question: str
    source_question: str
    source_bloom: str
    source_difficulty: str
    target_bloom: str
    target_difficulty: str
    predicted_bloom: str
    predicted_difficulty: str
    confidence: float
    attempts: int
    generation_time: float
    concept_match_score: str
    rejection_reason: str
    prompt_used: str
    explanation: str
    validation_status: str
    attempts_list: list = None
    prompt_tokens: int = 0
    completion_tokens: int = 0
    flan_calls: int = 0
    deberta_calls: int = 0
    candidates_generated: int = 0
    candidates_validated: int = 0
    candidates_rejected_before_validation: int = 0
    retry_history: list = None



REJECTION_LOGS = []


def run_warmup() -> bool:
    """Run a test inference pass on all three models.

    Returns True if every stage succeeded, False otherwise.
    Failures are printed with full tracebacks -- never silenced.
    """
    global _warmup_succeeded
    print("Running warm-up inference...")
    warmup_start = time.time()
    try:
        # Stage 1 -- DeBERTa
        _bloom, _diff, _conf = classify_text("What is database normalization?")
        print("  [warmup] DeBERTa OK  -> %s / %s / %.1f%%" % (_bloom, _diff, _conf))

        # Stage 2 -- FLAN-T5
        if flan_model is not None:
            meta_params = [n for n, p in flan_model.named_parameters() if p.device.type == "meta"]
            if meta_params:
                raise RuntimeError(
                    "FLAN-T5 has %d parameter(s) on 'meta' device after loading." % len(meta_params)
                )

            prompt_new = build_prompt(
                question="What is database normalization?",
                source_bloom="Remember",
                target_bloom="Remember",
                domain="CS",
                topic="Database",
            )
            inputs_new = flan_tokenizer(prompt_new, return_tensors="pt", truncation=True, max_length=128)
            with torch.inference_mode():
                out = flan_model.generate(**inputs_new, max_new_tokens=10, num_beams=1)
            decoded = flan_tokenizer.decode(out[0], skip_special_tokens=True)
            print("  [warmup] FLAN-T5 OK  -> '%s'" % decoded[:80])

        # Stage 3 -- SentenceTransformer
        st = _get_st_model()
        if st is not None:
            _ = st.encode(["warmup sentence a", "warmup sentence b"])
            print("  [warmup] SentenceTransformer OK")

        warmup_duration = time.time() - warmup_start
        print("Warm-up completed successfully in %.2f s." % warmup_duration)
        _warmup_succeeded = True
        return True

    except Exception as exc:
        import traceback as _tb
        print("[ERROR] Warm-up inference FAILED -- models are NOT ready for inference.")
        print(_tb.format_exc())
        _warmup_succeeded = False
        return False


def load_models():
    """Load DeBERTa, FLAN-T5, and SentenceTransformer sequentially.

    ROOT CAUSE: Concurrent from_pretrained() calls across multiple threads
    corrupt PyTorch's internal tensor allocation state in Transformers 5.x,
    leaving some FLAN-T5 parameters on the 'meta' device. Sequential loading
    inside a single thread is the correct and proven fix.
    """
    global deberta_model, deberta_tokenizer, flan_model, flan_tokenizer, _st_model, _load_errors

    # Fast-path: already loaded
    if (
        deberta_model is not None
        and deberta_tokenizer is not None
        and flan_model is not None
        and flan_tokenizer is not None
        and _st_model is not None
    ):
        print("[CACHE] All models already fully loaded.")
        return

    with MODEL_LOCK:
        # Double-check inside lock
        if (
            deberta_model is not None
            and deberta_tokenizer is not None
            and flan_model is not None
            and flan_tokenizer is not None
            and _st_model is not None
        ):
            print("[CACHE] All models already fully loaded (verified under lock).")
            return

        start_time = time.time()
        load_times = {}
        _load_errors = []

        print("\n" + "=" * 50)
        print("           MODEL LOADING (sequential)")
        print("=" * 50)

        # -- 1. DeBERTa -------------------------------------------------------
        t0 = time.time()
        try:
            if os.path.exists(DEBERTA_PATH):
                print("Loading DeBERTa...")
                deberta_tokenizer = AutoTokenizer.from_pretrained(DEBERTA_PATH)
                deberta_model = AutoModelForSequenceClassification.from_pretrained(
                    DEBERTA_PATH
                )
                device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
                deberta_model = deberta_model.to(device)
                deberta_model.eval()
                load_times["DeBERTa"] = time.time() - t0
                print("DeBERTa .............. OK (%.2f s) on %s" % (load_times["DeBERTa"], device))
            else:
                msg = "DeBERTa path not found: %s" % DEBERTA_PATH
                _load_errors.append(msg)
                print("[ERROR] %s" % msg)
        except Exception as exc:
            import traceback as _tb
            msg = "DeBERTa load failed: %s" % exc
            _load_errors.append(msg)
            print("[ERROR] %s" % msg)
            print(_tb.format_exc())

        # -- 2. FLAN-T5 -------------------------------------------------------
        t0 = time.time()
        try:
            if os.path.exists(FLAN_PATH):
                print("Loading FLAN-T5...")
                flan_tokenizer = AutoTokenizer.from_pretrained(FLAN_PATH)
                flan_model = AutoModelForSeq2SeqLM.from_pretrained(FLAN_PATH)
                device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
                flan_model = flan_model.to(device)
                flan_model.eval()
                # Verify no meta-device tensors
                meta_params = [
                    n for n, p in flan_model.named_parameters()
                    if p.device.type == "meta"
                ]
                if meta_params:
                    msg = (
                        "FLAN-T5 loaded with %d meta-device param(s). "
                        "Model is unusable. Clearing." % len(meta_params)
                    )
                    _load_errors.append(msg)
                    print("[ERROR] %s" % msg)
                    flan_model = None
                    flan_tokenizer = None
                else:
                    load_times["FLAN-T5"] = time.time() - t0
                    print("FLAN-T5 .............. OK (%.2f s)" % load_times["FLAN-T5"])
            else:
                msg = "FLAN-T5 path not found: %s" % FLAN_PATH
                _load_errors.append(msg)
                print("[ERROR] %s" % msg)
        except Exception as exc:
            import traceback as _tb
            msg = "FLAN-T5 load failed: %s" % exc
            _load_errors.append(msg)
            print("[ERROR] %s" % msg)
            print(_tb.format_exc())

        # -- 3. SentenceTransformer -------------------------------------------
        t0 = time.time()
        try:
            from sentence_transformers import SentenceTransformer
            print("Loading SentenceTransformer...")
            _st_model = SentenceTransformer("all-MiniLM-L6-v2")
            load_times["SentenceTransformer"] = time.time() - t0
            print("SentenceTransformer .. OK (%.2f s)" % load_times["SentenceTransformer"])
        except Exception as exc:
            import traceback as _tb
            msg = "SentenceTransformer load failed: %s" % exc
            _load_errors.append(msg)
            print("[ERROR] %s" % msg)
            print(_tb.format_exc())

        total_time = time.time() - start_time
        print("=" * 50)
        if _load_errors:
            print("[ERROR] %d model(s) failed to load: %s" % (len(_load_errors), _load_errors))
        else:
            print("All models loaded successfully.")
        print("Total load time: %.2f s" % total_time)
        print("=" * 50 + "\n")

        # -- Warm-up ----------------------------------------------------------
        if flan_model is not None and deberta_model is not None:
            run_warmup()
        else:
            print("[WARNING] Skipping warm-up: one or more core models failed to load.")

        mem_mb = get_current_process_memory()
        print("Memory after startup: %.2f MB" % mem_mb)

        # --- Default model confirmation banner ------------------------------
        print("\n" + "-" * 40)
        print("Default FLAN Model : flan_t5_model")
        print("Prompt Template    : Training Prompt")
        print("-" * 40 + "\n")


def sanitize_generated_question(text: str) -> str:
    """Post-process FLAN output to remove awkward instruction-style wording.

    Handles patterns like:
      'supplied X'  →  'X'
      'provided dataset'        →  'dataset'
      'given dataset'           →  'dataset'

    We do NOT automatically replace Bloom verbs (e.g. Execute -> Apply)
    as this alters the generated cognitive level and breaks classifier validation.
    """
    # 1. Remove the word "supplied" wherever it appears (keeps the noun after
    # it)
    text = re.sub(r"\bsupplied\s+", "", text, flags=re.IGNORECASE)

    # 2. Remove "provided" / "given" directly before
    # dataset/data/example/query/scenario
    text = re.sub(
        r"\b(provided|given)\s+(dataset|data|example|query|scenario)\b",
        r"\2",
        text,
        flags=re.IGNORECASE,
    )

    # 3. Collapse any double spaces introduced by removals and re-strip
    text = re.sub(r"  +", " ", text).strip()

    return text


# =====================================================================
# PREPROCESSING & EXPLANATION LAYER
# =====================================================================

# Pre-compiled regex patterns and module-level lookups for Preprocessing Layer
RE_CLEAN_1 = re.compile(r'^[\(\[][Qq]?\d+[\)\]]\s*')
RE_CLEAN_2 = re.compile(r'^(?:[Qq](?:uestion|uest|uery)?\s*[\.\-\:]?\s*)?\d+(?:[\.\)\:]\s*|\-\s+|\s+)')

RE_NON_ALPHANUM = re.compile(r'[^a-zA-Z0-9]')

_CONCEPT_MAPPINGS = {
    "normalization in dbms": "Database Normalization",
    "normalization in database management systems": "Database Normalization",
    "normalization": "Database Normalization",
    "database normalization": "Database Normalization",
    "un-normalised relational table": "Database Normalization",
    "un normalised relational table": "Database Normalization",
    "acid properties of a database transaction": "ACID Properties",
    "acid properties": "ACID Properties",
    "acid": "ACID Properties",
    "osi model": "OSI Model",
    "osi": "OSI Model",
    "tcp congestion control": "TCP Congestion Control",
    "tcp": "TCP Protocol",
    "udp": "UDP Protocol",
    "binary search tree": "Binary Search Tree",
    "avl tree": "AVL Tree",
    "binary search": "Binary Search",
    "big data": "Big Data",
    "deadlock": "Deadlock",
    "deadlock avoidance": "Deadlock Avoidance",
    "deadlock detection": "Deadlock Detection",
    "page fault": "Page Fault",
    "cpu scheduling": "CPU Scheduling",
    "agile methodology": "Agile Methodology",
    "software design principles": "Software Design Principles",
    "solid principles": "SOLID Principles",
    "design patterns": "Design Patterns",
    "version control": "Version Control",
    "graph traversal": "Graph Traversal",
    "hash tables": "Hash Tables",
    "heaps": "Heaps",
    "gradient descent": "Gradient Descent",
    "learning paradigms": "Learning Paradigms",
    "convolutional neural networks": "Convolutional Neural Networks",
    "search algorithms": "Search Algorithms",
}

_ACRONYMS = {
    "dbms": "DBMS", "tcp": "TCP", "udp": "UDP", "osi": "OSI", "sql": "SQL",
    "nosql": "NoSQL", "lru": "LRU", "fifo": "FIFO", "acid": "ACID", "bcnf": "BCNF",
    "3nf": "3NF", "2nf": "2NF", "1nf": "1NF", "uml": "UML", "ip": "IP",
    "dhcp": "DHCP", "arp": "ARP", "avl": "AVL", "cidr": "CIDR", "rsa": "RSA",
    "cpu": "CPU", "solid": "SOLID", "gui": "GUI", "api": "API", "dns": "DNS",
    "lan": "LAN", "wan": "WAN", "mac": "MAC", "http": "HTTP", "https": "HTTPS",
    "xml": "XML", "json": "JSON",
}


@functools.lru_cache(maxsize=1024)
def clean_source_question(question: str) -> str:
    """
    Remove leading question numbering, prefixes (e.g. Q1., Q2), (3), 1., 2)
    and extra whitespace. Do not modify academic content.
    """
    if not question:
        return ""
    q = question.strip()
    q_clean = RE_CLEAN_1.sub('', q)
    q_clean = RE_CLEAN_2.sub('', q_clean)
    if not q_clean.strip():
        return q
    return q_clean.strip()


@functools.lru_cache(maxsize=1024)
def format_academic_concept(concept: str) -> str:
    """
    Formats the extracted concept to be a clean, Title-Cased academic name,
    preserving technical acronyms (DBMS, TCP, UDP, OSI, SQL, NoSQL, LRU, FIFO, ACID, BCNF, 3NF, etc.)
    and mapping specific messy strings to proper academic names.
    """
    if not concept:
        return "Technical Concept"
        
    concept_lower = concept.lower().strip()
    
    if concept_lower in _CONCEPT_MAPPINGS:
        return _CONCEPT_MAPPINGS[concept_lower]
        
    if "normalization" in concept_lower and ("dbms" in concept_lower or "database" in concept_lower):
        return "Database Normalization"
        
    words = concept_lower.split()
    formatted_words = []
    for w in words:
        w_clean = RE_NON_ALPHANUM.sub('', w)
        if w_clean in _ACRONYMS:
            val = _ACRONYMS[w_clean]
            formatted_word = w.replace(w_clean, val)
            formatted_words.append(formatted_word)
        else:
            formatted_words.append(w.capitalize())
            
    res = " ".join(formatted_words).strip()
    
    res_words = res.split()
    if len(res_words) > 5:
        res = " ".join(res_words[:3])
        
    if not res:
        return "Technical Concept"
        
    return res


@functools.lru_cache(maxsize=1024)
def normalize_academic_concept(question: str) -> str:
    """
    Extracts the core phrase from a cleaned question and formats it
    as a clean Title-Cased academic name.
    """
    raw_concept = extract_core_phrase(question)
    return format_academic_concept(raw_concept)


def generate_dynamic_explanation(
    question: str, required_concept: str, bloom: str, difficulty: str
) -> str:
    """
    Generates a concise, 2-3 sentence academic explanation for a question's classification,
    using the formatted academic concept.
    """
    # Use pre-formatted concept if available, avoiding duplicate normalize/format passes
    if required_concept:
        concept = required_concept
    else:
        concept = normalize_academic_concept(question)

    # Standardize difficulty label for output explanation
    diff_word = (
        "Medium"
        if difficulty == "Moderate"
        else ("Hard" if difficulty == "Difficult" else difficulty)
    )

    templates = {
        "Remember": "This question asks the learner to recall fundamental knowledge about {concept}.",
        "Understand": "This question requires explaining the principles of {concept}.",
        "Apply": "This question requires applying {concept} to solve a practical problem.",
        "Analyze": "This question requires examining the components or relationships within {concept}.",
        "Evaluate": "This question requires assessing or justifying decisions involving {concept}.",
        "Create": "This question requires designing or constructing a solution using {concept}."
    }

    s1 = templates.get(bloom, "This question focuses on the principles of {concept}.").format(concept=concept)
    s2 = f"A {bloom} level classification is appropriate because of the cognitive load required to address the topic."
    s3 = f"The difficulty is {diff_word} to match the depth of reasoning expected in the solution."

    return f"{s1} {s2} {s3}"


def classify_text(text: str, return_details: bool = False):
    if deberta_model is None or deberta_tokenizer is None:
        raise RuntimeError("DeBERTa model is not loaded.")

    try:
        t_tok_start = time.perf_counter()
        inputs = deberta_tokenizer(
            text, return_tensors="pt", truncation=True, max_length=512
        )
        t_tok = (time.perf_counter() - t_tok_start) * 1000.0

        t_inf_start = time.perf_counter()
        device = next(deberta_model.parameters()).device
        inputs = {k: v.to(device) for k, v in inputs.items()}
        with torch.inference_mode():
            outputs = deberta_model(**inputs)
            logits = outputs.logits
        t_inf = (time.perf_counter() - t_inf_start) * 1000.0

        t_bloom_start = time.perf_counter()
        predicted_class_id = logits.argmax().item()
        index_to_bloom = {
            0: "Remember",
            1: "Understand",
            2: "Apply",
            3: "Analyze",
            4: "Evaluate",
            5: "Create",
        }
        bloom = index_to_bloom.get(predicted_class_id, "Analyze")
        t_bloom = (time.perf_counter() - t_bloom_start) * 1000.0

        t_diff_start = time.perf_counter()
        diff = BLOOM_TO_DIFFICULTY.get(bloom, "Medium")
        t_diff = (time.perf_counter() - t_diff_start) * 1000.0

        t_conf_start = time.perf_counter()
        probs = torch.nn.functional.softmax(logits, dim=-1)
        confidence = probs[0][predicted_class_id].item() * 100
        t_conf = (time.perf_counter() - t_conf_start) * 1000.0

        if return_details:
            return bloom, diff, round(confidence, 2), {
                "tokenizer": t_tok,
                "deberta_inference": t_inf,
                "bloom_mapping": t_bloom,
                "difficulty_mapping": t_diff,
                "confidence_calc": t_conf,
            }

        return bloom, diff, round(confidence, 2)
    except Exception as e:
        print(f"Error in classification: {e}")
        if return_details:
            return "Unknown", "Unknown", 0.0, {
                "tokenizer": 0.0,
                "deberta_inference": 0.0,
                "bloom_mapping": 0.0,
                "difficulty_mapping": 0.0,
                "confidence_calc": 0.0,
            }
        return "Unknown", "Unknown", 0.0



def generate_single_variant(
    question,
    source_bloom,
    source_difficulty,
    target_bloom,
    target_difficulty,
    domain="General Computer Science",
    topic="General",
    required_concept="",
    attempt_number=1,
    config_mode=None,
    mode_name=None,
    failure_hint="",
    return_timings=False,
):
    """Generate variant question candidates using the FLAN-T5 model (Mode E)."""
    if config_mode is None:
        config_mode = MODES[ACTIVE_MODE]
        mode_name = ACTIVE_MODE
    if mode_name is None:
        mode_name = ACTIVE_MODE

    if flan_model is None or flan_tokenizer is None:
        raise RuntimeError("FLAN-T5 model is not loaded.")

    topic_to_use = required_concept if required_concept else topic

    t_prompt_start = time.perf_counter()
    prompt = build_prompt(
        question=question,
        source_bloom=source_bloom,
        target_bloom=target_bloom,
        domain=domain,
        topic=topic_to_use,
    )

    # Adaptive Retry Prompt Scaling
    if attempt_number == 2:
        # Attempt 2: Training prompt + Minimal retry guidance (e.g., failure hint)
        if failure_hint:
            prompt += f"\nPrevious attempt failed because: {failure_hint}. Avoid this."
    elif attempt_number == 3:
        # Attempt 3: Training prompt + Retry guidance + concise concept reminder
        if failure_hint:
            prompt += f"\nPrevious attempt failed because: {failure_hint}. Avoid this."
        if topic_to_use:
            # Issue 10: Single concise concept reminder replaces four redundant lines.
            prompt += f"\nOriginal Concept: {topic_to_use}\nThe generated question must explicitly preserve this concept."
    elif attempt_number >= 4:
        # Attempt 4+: Training prompt + Retry guidance + Concept reminder + Bloom verb reminder
        if failure_hint:
            prompt += f"\nPrevious attempt failed because: {failure_hint}. Avoid this."
        if topic_to_use:
            # Issue 10: Single concise concept reminder.
            prompt += f"\nOriginal Concept: {topic_to_use}\nThe generated question must explicitly preserve this concept."
        # Mode E Bloom-level verb guidance
        mode_e_guidance = {
            "Remember": "\nStart: Define, Identify, List, Name, State, Recall. Forbid: Explain, Discuss, Analyze, Evaluate, Design.",
            "Understand": "\nStart: Explain, Describe, Discuss, Summarize, Interpret. Forbid: Define, List, Apply, Analyze, Design.",
            "Apply": "\nStart: Apply, Implement, Use, Demonstrate, Solve, Calculate. Forbid: Define, Explain, Analyze, Evaluate, Design.",
            "Analyze": "\nStart: Analyze, Compare, Differentiate, Examine, Contrast, Distinguish. Forbid: Define, Explain, Apply, Evaluate, Design, suitable, best, recommend, judge.",
            "Evaluate": "\nStart: Evaluate, Assess, Critique, Justify, Defend, Determine. Forbid: Design, Create, Develop, Formulate, Construct.",
            "Create": "\nStart: Design, Create, Develop, Formulate, Construct. Forbid: Define, Explain, Apply, Analyze, Evaluate, Describe."
        }
        guidance = mode_e_guidance.get(target_bloom, "")
        if guidance:
            prompt += guidance
    t_prompt_ms = (time.perf_counter() - t_prompt_start) * 1000.0

    t_flan_start = time.perf_counter()
    inputs = flan_tokenizer(prompt, return_tensors="pt", truncation=True, max_length=512)
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    inputs = {k: v.to(device) for k, v in inputs.items()}

    # Resolve candidate sequence counts based on Bloom level
    num_seq = NUM_CANDIDATES_BY_BLOOM.get(target_bloom, 3)

    gen_kwargs = {
        "max_new_tokens": config_mode["max_new_tokens"],
        "repetition_penalty": config_mode["repetition_penalty"],
        "no_repeat_ngram_size": config_mode["no_repeat_ngram_size"],
        "num_return_sequences": num_seq,
        "do_sample": True,
        "temperature": config_mode.get("temperature", 0.7),
        "top_p": config_mode.get("top_p", 0.95),
    }

    with torch.inference_mode():
        outputs = flan_model.generate(**inputs, **gen_kwargs)

    result = [flan_tokenizer.decode(out, skip_special_tokens=True).strip() for out in outputs]
    t_flan_ms = (time.perf_counter() - t_flan_start) * 1000.0

    if return_timings:
        return result, prompt, t_prompt_ms, t_flan_ms
    return result, prompt


def calculate_validation_status(
    predicted_bloom, predicted_diff, target_bloom, target_diff
):
    """Return Exact Match, Adjacent Match, or Mismatch based on Bloom level."""
    if predicted_bloom == target_bloom:
        return "Exact Match"
    bloom_levels = ["Remember", "Understand", "Apply", "Analyze", "Evaluate", "Create"]
    try:
        p_idx = bloom_levels.index(predicted_bloom)
        t_idx = bloom_levels.index(target_bloom)
        if abs(p_idx - t_idx) == 1:
            return "Adjacent Match"
    except ValueError:
        pass
    return "Mismatch"



def extract_questions_from_file(file):
    filename = secure_filename(file.filename).lower()
    questions = []

    if filename.endswith(".csv"):
        stream = io.StringIO(file.stream.read().decode("utf8", errors="ignore"))
        reader = csv.reader(stream)
        headers = next(reader, None)
        q_idx = 0
        if headers:
            for i, h in enumerate(headers):
                if "question" in h.lower():
                    q_idx = i
                    break
        for row in reader:
            if len(row) > q_idx and row[q_idx].strip():
                questions.append(row[q_idx].strip())

    elif filename.endswith(".txt"):
        content = file.read().decode("utf-8", errors="ignore")
        # Split by double/multiple newlines (possibly with whitespace on the blank lines)
        blocks = re.split(r'\r?\n\s*\r?\n', content)
        for block in blocks:
            if block.strip():
                questions.append(block.strip())

    elif filename.endswith(".pdf"):
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        for line in text.splitlines():
            # simple heuristic for lines that might be questions
            if line.strip() and len(line.strip()) > 10:
                questions.append(line.strip())

    elif filename.endswith(".docx"):
        doc = docx.Document(file)
        for para in doc.paragraphs:
            if para.text.strip():
                questions.append(para.text.strip())

    return questions


def update_batch_history_stats(session_id):
    if session_id not in MEMORY_STORE:
        return
    store = MEMORY_STORE[session_id]
    for hist in BATCH_HISTORY:
        if hist["session_id"] == session_id:
            hist["total_questions"] = store["total_questions"]
            hist["completed_questions"] = store["completed_questions"]
            
            easy, medium, hard = 0, 0, 0
            bloom_counts = { "Remember": 0, "Understand": 0, "Apply": 0, "Analyze": 0, "Evaluate": 0, "Create": 0 }
            total_conf = 0
            for r in store["results"]:
                d = r.get("difficulty")
                if d == 'Easy': easy += 1
                elif d in ['Medium', 'Moderate']: medium += 1
                elif d in ['Hard', 'Difficult']: hard += 1
                
                bl = r.get("bloom_level")
                if bl in bloom_counts:
                    bloom_counts[bl] += 1
                
                total_conf += float(r.get("confidence", 0))
            
            hist["easy_count"] = easy
            hist["medium_count"] = medium
            hist["hard_count"] = hard
            hist["bloom_counts"] = bloom_counts
            hist["avg_confidence"] = round((total_conf / len(store["results"])), 1) if len(store["results"]) > 0 else 0
            break


def process_batch_background(session_id):
    store = MEMORY_STORE.get(session_id)
    if not store:
        return

    store["status"] = "PROCESSING"
    store["start_time"] = time.time()
    batch_start_perf = time.perf_counter()

    questions = store["questions"]
    total_q = len(questions)

    import spacy_utils

    # --- Log One-Time Costs at Batch Start ---
    print(f"\n==================================================")
    print(f"Batch Started : Session {session_id} ({total_q} questions)")
    print(f"==================================================")
    print(f"Tokenizer Status          : Loaded ({DEBERTA_PATH})")
    print(f"SentenceTransformer Status : {'Loaded' if _st_model is not None else 'Unloaded / Lazy'}")
    print(f"spaCy Model Status        : {'Loaded' if spacy_utils._nlp is not None else 'Not Loaded (Will Lazy Load)'}")
    print(f"Cache Initialization      : Embedding Cache ({len(EMBEDDING_CACHE)} keys, Hits: {EMBEDDING_CACHE_HITS}, Misses: {EMBEDDING_CACHE_MISSES})")
    print(f"spaCy Doc Cache           : {len(spacy_utils._doc_cache)} documents")
    print(f"Thread Info               : ID {threading.get_ident()} ({threading.current_thread().name})")
    print(f"Memory Usage              : {get_current_process_memory():.2f} MB")
    print(f"==================================================\n")

    batch_timings = []

    try:
        for i, raw_q in enumerate(questions):
            if store["stop"]:
                store["status"] = "STOPPED"
                break

            q_num = i + 1

            # 1. Question Extraction
            t_ext_start = time.perf_counter()
            q = raw_q
            t_ext = (time.perf_counter() - t_ext_start) * 1000.0

            # 2. Input Validation
            t_val_start = time.perf_counter()
            is_valid = isinstance(q, str) and len(q.strip()) > 0
            t_val = (time.perf_counter() - t_val_start) * 1000.0

            # 3. Text Cleaning
            t_clean_start = time.perf_counter()
            cleaned_q = clean_source_question(q)
            t_clean = (time.perf_counter() - t_clean_start) * 1000.0

            # 4. spaCy Processing (Deferred to Phase 2 after classification)
            t_spacy = 0.0

            # 5. SentenceTransformer (Deferred to Phase 2 after classification)
            t_st = 0.0

            # 6. Cache Lookup
            t_cache = 0.0

            # 7-11. Tokenizer, DeBERTa Inference, Bloom Mapping, Difficulty Mapping, Confidence Calculation
            bloom, diff, conf, det_timings = classify_text(cleaned_q, return_details=True)
            t_tok = det_timings["tokenizer"]
            t_inf = det_timings["deberta_inference"]
            t_bloom = det_timings["bloom_mapping"]
            t_diff = det_timings["difficulty_mapping"]
            t_conf = det_timings["confidence_calc"]

            # 12. Explanation Generation
            t_exp_start = time.perf_counter()
            required_concept = normalize_academic_concept(cleaned_q)
            explanation = generate_dynamic_explanation(cleaned_q, required_concept, bloom, diff)
            t_exp = (time.perf_counter() - t_exp_start) * 1000.0

            timestamp_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            result = {
                "id": q_num,
                "question": cleaned_q,
                "original_question": q,
                "bloom_level": bloom,
                "difficulty": diff,
                "confidence": conf,
                "explanation": explanation,
                "timestamp": timestamp_str,
                "status": (
                    "Verified"
                    if float(conf) >= 95
                    else ("Review" if float(conf) >= 80 else "Low Confidence")
                ),
                "previous_classification": None,
                "variants": [],
                "notes": "",
                "tags": [],
            }

            # 13. Session Update
            t_sess_start = time.perf_counter()
            store["completed_questions"] = q_num
            t_sess = (time.perf_counter() - t_sess_start) * 1000.0

            # 14. Result Storage
            t_store_start = time.perf_counter()
            store["results"].append(result)
            t_store = (time.perf_counter() - t_store_start) * 1000.0

            stages = {
                "Question Extraction": t_ext,
                "Input Validation": t_val,
                "Text Cleaning": t_clean,
                "spaCy Processing": t_spacy,
                "SentenceTransformer": t_st,
                "Cache Lookup": t_cache,
                "Tokenizer": t_tok,
                "DeBERTa Inference": t_inf,
                "Bloom Mapping": t_bloom,
                "Difficulty Mapping": t_diff,
                "Confidence Calculation": t_conf,
                "Explanation Generation": t_exp,
                "Session Update": t_sess,
                "Result Storage": t_store,
            }

            total_q_ms = sum(stages.values())
            slowest_stage_name, slowest_stage_val = max(stages.items(), key=lambda x: x[1])

            batch_timings.append({
                "q_num": q_num,
                "total_ms": total_q_ms,
                "stages": stages,
                "slowest_stage": (slowest_stage_name, slowest_stage_val),
            })

            # Print per-question instrumentation log (if enabled)
            if ENABLE_PERFORMANCE_LOGS:
                print(f"==================================================")
                print(f"Question {q_num}")
                print(f"==================================================")
                for s_name, s_val in stages.items():
                    print(f"{s_name:<25}: {s_val:8.2f} ms")
                print(f"\nTOTAL                    : {total_q_ms:8.2f} ms")
                print(f"Slowest Stage            : {slowest_stage_name} ({slowest_stage_val:.2f} ms)")
                print(f"==================================================\n")

        if not store["stop"] and store["completed_questions"] == store["total_questions"]:
            store["status"] = "COMPLETED"

    except Exception as e:
        import traceback
        print("Exception in background batch processing:")
        traceback.print_exc()
        store["status"] = "FAILED"

    # --- Batch Performance Report & Group Comparisons ---
    if batch_timings:
        total_batch_ms = (time.perf_counter() - batch_start_perf) * 1000.0
        n = len(batch_timings)
        avg_q_ms = total_batch_ms / n

        fastest = min(batch_timings, key=lambda x: x["total_ms"])
        slowest = max(batch_timings, key=lambda x: x["total_ms"])
        sorted_by_slow = sorted(batch_timings, key=lambda x: x["total_ms"], reverse=True)
        top5 = sorted_by_slow[:5]

        # Group comparison (Q1-10 vs Q11+)
        q1_10 = batch_timings[:10]
        q11_plus = batch_timings[10:]

        avg_q1_10 = sum(x["total_ms"] for x in q1_10) / len(q1_10) if q1_10 else 0.0
        avg_q11_plus = sum(x["total_ms"] for x in q11_plus) / len(q11_plus) if q11_plus else 0.0
        diff_groups = avg_q1_10 - avg_q11_plus if q11_plus else 0.0

        # Stage averages across batch
        stage_totals = {k: 0.0 for k in batch_timings[0]["stages"].keys()}
        stage_q1_10_totals = {k: 0.0 for k in batch_timings[0]["stages"].keys()}
        stage_q11_plus_totals = {k: 0.0 for k in batch_timings[0]["stages"].keys()}

        for item in batch_timings:
            for s_k, s_v in item["stages"].items():
                stage_totals[s_k] += s_v

        for item in q1_10:
            for s_k, s_v in item["stages"].items():
                stage_q1_10_totals[s_k] += s_v

        for item in q11_plus:
            for s_k, s_v in item["stages"].items():
                stage_q11_plus_totals[s_k] += s_v

        stage_avgs = {k: v / n for k, v in stage_totals.items()}
        stage_q1_10_avgs = {k: v / len(q1_10) for k, v in stage_q1_10_totals.items()} if q1_10 else {}
        stage_q11_plus_avgs = {k: v / len(q11_plus) for k, v in stage_q11_plus_totals.items()} if q11_plus else {}

        largest_bottleneck_name, largest_bottleneck_val = max(stage_avgs.items(), key=lambda x: x[1])
        bottleneck_pct = (largest_bottleneck_val / (sum(stage_avgs.values()) or 1.0)) * 100.0

        # Warm-up Bottlenecks
        warmup_bottlenecks = []
        if q11_plus:
            for s_k in stage_avgs.keys():
                q1_10_avg = stage_q1_10_avgs.get(s_k, 0.0)
                q11_avg = stage_q11_plus_avgs.get(s_k, 0.0)
                if q1_10_avg > 2.0 * (q11_avg or 1.0) and q1_10_avg > 20.0:
                    warmup_bottlenecks.append((s_k, q1_10_avg, q11_avg))

        if ENABLE_PERFORMANCE_LOGS:
            print(f"\n==================================================")
            print(f"            Group Performance Analysis            ")
            print(f"==================================================")
            print(f"Questions 1-10 Average Classification Time  : {avg_q1_10:.2f} ms")
            if q11_plus:
                print(f"Questions 11-{n} Average Classification Time: {avg_q11_plus:.2f} ms")
                print(f"Difference (First 10 vs Later Questions)  : {diff_groups:+.2f} ms")
            print(f"\nAverage Stage Times Across Batch:")
            for s_k, s_v in stage_avgs.items():
                q1_avg = stage_q1_10_avgs.get(s_k, 0.0)
                q11_avg = stage_q11_plus_avgs.get(s_k, 0.0) if q11_plus else 0.0
                print(f"  {s_k:<25}: Batch Avg {s_v:7.2f} ms | Q1-10 {q1_avg:7.2f} ms | Q11+ {q11_avg:7.2f} ms")

            if warmup_bottlenecks:
                print(f"\n[WARM-UP BOTTLENECKS DETECTED]")
                for wb_name, wb_q1, wb_q11 in warmup_bottlenecks:
                    print(f"  - Possible Warm-up Bottleneck: {wb_name} (Q1-10 Avg: {wb_q1:.2f} ms -> Q11+ Avg: {wb_q11:.2f} ms)")

            print(f"\n==================================================")
            print(f"            Batch Performance Report              ")
            print(f"==================================================")
            print(f"Total Questions         : {n}")
            print(f"Total Time              : {total_batch_ms:.2f} ms")
            print(f"Average Time            : {avg_q_ms:.2f} ms")
            print(f"Fastest Question        : Question {fastest['q_num']} ({fastest['total_ms']:.2f} ms)")
            print(f"Slowest Question        : Question {slowest['q_num']} ({slowest['total_ms']:.2f} ms)")
            print(f"\nTop 5 Slowest Questions:")
            for idx_top, item in enumerate(top5, 1):
                s_name, s_val = item['slowest_stage']
                print(f"  {idx_top}. Question {item['q_num']} - {item['total_ms']:.2f} ms (Slowest Stage: {s_name} {s_val:.2f} ms)")

            print(f"\nAverage Time Per Stage:")
            for s_name, s_val in stage_avgs.items():
                print(f"  {s_name:<25}: {s_val:8.2f} ms")

            print(f"\nLargest Bottleneck      : {largest_bottleneck_name} ({largest_bottleneck_val:.2f} ms / {bottleneck_pct:.1f}% of per-question latency)")

            print(f"\nRecommendations:")
            if warmup_bottlenecks:
                top_wb = warmup_bottlenecks[0]
                print(f"  - Pre-warm '{top_wb[0]}' during system startup (load_models/run_warmup) to eliminate the initial {top_wb[1]:.0f} ms spike.")
            else:
                print(f"  - System stages exhibit consistent execution across batch questions.")
            print(f"==================================================\n")

    print(f"Batch {store['status']}: Session {session_id} ({store['completed_questions']}/{store['total_questions']} completed)")

    # Update History
    for hist in BATCH_HISTORY:
        if hist["session_id"] == session_id:
            hist["status"] = store["status"]
            hist["completed_questions"] = store["completed_questions"]
            break

    update_batch_history_stats(session_id)




@app.route("/")
@app.route("/landing")
def landing():
    return render_template("landing.html")


@app.route("/dashboard")
def index():
    return render_template("index.html")


@app.route("/auth")
def auth():
    return render_template("auth.html")



@app.route("/health", methods=["GET"])
def health():
    """Production health check.

    models_ready is True only when:
    - all model objects are non-None
    - no parameters are on the meta device
    - warm-up inference passed
    """
    deberta_ok = deberta_model is not None and deberta_tokenizer is not None
    flan_ok = flan_model is not None and flan_tokenizer is not None
    st_ok = _st_model is not None

    # Verify FLAN-T5 has no meta-device tensors
    flan_meta_free = False
    if flan_ok:
        flan_meta_free = not any(
            p.device.type == "meta" for p in flan_model.parameters()
        )
    else:
        flan_meta_free = False

    models_ready = deberta_ok and flan_ok and st_ok and flan_meta_free and _warmup_succeeded

    reasons = []
    if not deberta_ok:
        reasons.append("DeBERTa not loaded")
    if not flan_ok:
        reasons.append("FLAN-T5 not loaded")
    if flan_ok and not flan_meta_free:
        reasons.append("FLAN-T5 has meta-device parameters (parallel load race)")
    if not st_ok:
        reasons.append("SentenceTransformer not loaded")
    if not _warmup_succeeded:
        reasons.append("warm-up inference has not completed successfully")
    if _load_errors:
        reasons.extend(_load_errors)

    deberta_device = "cpu"
    if deberta_model is not None:
        try:
            deberta_device = str(next(deberta_model.parameters()).device)
        except Exception:
            pass

    flan_device = "cpu"
    if flan_model is not None:
        try:
            flan_device = str(next(flan_model.parameters()).device)
        except Exception:
            pass

    st_device = "cpu"
    if _st_model is not None:
        try:
            st_device = str(_st_model.device)
        except Exception:
            try:
                st_device = str(next(_st_model.parameters()).device)
            except Exception:
                pass

    return jsonify({
        "status": "healthy" if models_ready else "initializing",
        "models_ready": models_ready,
        "warmup_succeeded": _warmup_succeeded,
        "uptime_seconds": round(time.time() - STARTUP_TIMESTAMP, 2),
        "active_mode": ACTIVE_MODE,
        "model_versions": {
            "deberta": DEBERTA_PATH,
            "flan": FLAN_PATH,
            "sentence_transformer": "all-MiniLM-L6-v2",
        },
        "startup_timestamp": STARTUP_TIMESTAMP,
        "reasons": reasons,
        "load_errors": _load_errors,
        "models": {
            "deberta": {
                "loaded": deberta_ok,
                "device": deberta_device,
                "warmup": _warmup_succeeded
            },
            "flan_t5": {
                "loaded": flan_ok,
                "device": flan_device,
                "warmup": _warmup_succeeded
            },
            "sentence_transformer": {
                "loaded": st_ok,
                "device": st_device
            }
        },
        "validation_engine": {
            "loaded": True,
            "pipeline": "7-stage"
        },
        "embedding_cache": {
            "hits": EMBEDDING_CACHE_HITS,
            "misses": EMBEDDING_CACHE_MISSES
        },
        "memory_usage_mb": get_current_process_memory(),
        "python_version": platform.python_version(),
        "torch_version": torch.__version__,
        "transformers_version": transformers.__version__,
        "spacy_model": "en_core_web_sm"
    })


@app.route("/classify", methods=["POST"])
def classify():
    t_total_start = time.perf_counter()
    data = request.json or {}
    question = data.get("question", "")
    if not question:
        return jsonify({"error": "Question is required"}), 400

    # 1. Text Cleaning
    t_clean_start = time.perf_counter()
    cleaned_q = clean_source_question(question)
    t_clean = (time.perf_counter() - t_clean_start) * 1000.0

    # 2. DeBERTa Classification
    bloom, diff, conf, det_timings = classify_text(cleaned_q, return_details=True)
    t_tok = det_timings["tokenizer"]
    t_inf = det_timings["deberta_inference"]
    t_bloom = det_timings["bloom_mapping"]
    t_diff = det_timings["difficulty_mapping"]
    t_conf = det_timings["confidence_calc"]

    # 3. Concept Normalization (extract_core_phrase & format_academic_concept)
    t_nac_start = time.perf_counter()

    t_ecp_start = time.perf_counter()
    raw_concept = extract_core_phrase(cleaned_q)
    t_ecp = (time.perf_counter() - t_ecp_start) * 1000.0

    t_fac_start = time.perf_counter()
    required_concept = format_academic_concept(raw_concept)
    t_fac = (time.perf_counter() - t_fac_start) * 1000.0

    t_nac = (time.perf_counter() - t_nac_start) * 1000.0

    # 4. Explanation Generation
    t_exp_start = time.perf_counter()
    explanation = generate_dynamic_explanation(cleaned_q, required_concept, bloom, diff)
    t_exp = (time.perf_counter() - t_exp_start) * 1000.0

    # 5. JSON Response & History
    t_json_start = time.perf_counter()
    new_classification = {
        "id": "manual-" + str(int(time.time() * 1000)),
        "question": cleaned_q,
        "bloom_level": bloom,
        "difficulty": diff,
        "confidence": conf,
        "explanation": explanation,
        "upload_time": datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%S.%fZ")[:-4] + "Z",
        "status": "Completed"
    }
    SESSION_STATE["manual_classifications"].append(new_classification)

    resp = jsonify(
        {
            "question": cleaned_q,
            "original_question": question,
            "bloom_level": bloom,
            "difficulty": diff,
            "confidence": conf,
            "explanation": explanation,
        }
    )
    t_json = (time.perf_counter() - t_json_start) * 1000.0
    t_total = (time.perf_counter() - t_total_start) * 1000.0

    print("-----------------------------------------", flush=True)
    print("CLASSIFY PERFORMANCE", flush=True)
    print("-----------------------------------------", flush=True)
    print(f"Cleaning                 : {t_clean:.2f} ms", flush=True)
    print(f"Tokenizer                : {t_tok:.2f} ms", flush=True)
    print(f"DeBERTa                  : {t_inf:.2f} ms", flush=True)
    print(f"Bloom Mapping            : {t_bloom:.2f} ms", flush=True)
    print(f"Difficulty Mapping       : {t_diff:.2f} ms", flush=True)
    print(f"Confidence               : {t_conf:.2f} ms", flush=True)
    print(f"normalize_academic_concept : {t_nac:.2f} ms", flush=True)
    print(f"extract_core_phrase      : {t_ecp:.2f} ms", flush=True)
    print(f"format_academic_concept  : {t_fac:.2f} ms", flush=True)
    print(f"Explanation Generation   : {t_exp:.2f} ms", flush=True)
    print(f"JSON Response            : {t_json:.2f} ms", flush=True)
    print(f"TOTAL                    : {t_total:.2f} ms", flush=True)
    print("-----------------------------------------", flush=True)

    return resp


@app.route("/manual-classifications", methods=["GET"])
def get_manual_classifications():
    return jsonify(SESSION_STATE["manual_classifications"])



@app.route("/parse-upload", methods=["POST"])
def parse_upload():
    if "file" not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No selected file"}), 400

    try:
        filename = secure_filename(file.filename)
        questions = extract_questions_from_file(file)
        questions = questions[:1000]  # Limit to 1000 for practical reasons

        if not questions:
            return (
                jsonify(
                    {
                        "error": "No questions extracted from file. Ensure file is not empty."
                    }
                ),
                400,
            )

        return jsonify({"filename": filename, "questions": questions})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/upload-batch", methods=["POST"])
def upload_batch():
    t_start = time.perf_counter()
    data = request.json
    if not data or "questions" not in data:
        return jsonify({"error": "Missing questions data"}), 400

    questions = data.get("questions", [])
    filename = data.get("filename", "Batch Upload")

    if not questions:
        return jsonify({"error": "No questions to process."}), 400

    try:
        session_id = str(uuid.uuid4())
        MEMORY_STORE[session_id] = {
            "session_id": session_id,
            "filename": filename,
            "total_questions": len(questions),
            "questions": questions,
            "results": [],
            "status": "PENDING",
            "stop": False,
            "start_time": None,
            "completed_questions": 0,
        }
        session["batch_id"] = session_id

        BATCH_HISTORY.insert(
            0,
            {
                "session_id": session_id,
                "filename": filename,
                "upload_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "total_questions": len(questions),
                "completed_questions": 0,
                "status": "PENDING",
                "easy_count": 0,
                "medium_count": 0,
                "hard_count": 0,
                "bloom_counts": { "Remember": 0, "Understand": 0, "Apply": 0, "Analyze": 0, "Evaluate": 0, "Create": 0 },
                "avg_confidence": 0
            },
        )
        if len(BATCH_HISTORY) > 50:
            BATCH_HISTORY.pop()

        t_proc_ms = (time.perf_counter() - t_start) * 1000.0
        print(f"[PERF-API] POST /upload-batch Processing Time: {t_proc_ms:.2f} ms ({len(questions)} questions)")

        return jsonify(
            {
                "session_id": session_id,
                "total_questions": len(questions),
                "status": "PENDING",
            }
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/start-batch/<session_id>", methods=["POST"])
def start_batch(session_id):
    if session_id not in MEMORY_STORE:
        return jsonify({"error": "Invalid session ID"}), 404

    store = MEMORY_STORE[session_id]
    if store["status"] != "PENDING":
        return jsonify({"error": "Batch already started or completed"}), 400

    thread = threading.Thread(target=process_batch_background, args=(session_id,))
    thread.start()
    return jsonify({"message": "Processing started"})


@app.route("/batch-status/<session_id>", methods=["GET"])
def get_batch_status(session_id):
    t_route_start = time.perf_counter()
    if session_id not in MEMORY_STORE:
        return jsonify({"error": "Invalid session ID"}), 404

    t_proc_start = time.perf_counter()
    store = MEMORY_STORE[session_id]

    speed = 0.0
    eta = 0
    if store["start_time"] and store["completed_questions"] > 0:
        elapsed = time.time() - store["start_time"]
        if elapsed > 0:
            speed = store["completed_questions"] / elapsed
            remaining_q = store["total_questions"] - store["completed_questions"]
            if speed > 0:
                eta = int(remaining_q / speed)
    t_proc_ms = (time.perf_counter() - t_proc_start) * 1000.0

    t_json_start = time.perf_counter()
    response_obj = {
        "status": store["status"],
        "total": store["total_questions"],
        "processed": store["completed_questions"],
        "results": store["results"],
        "speed": round(speed, 1),
        "eta": eta,
    }
    json_bytes = json.dumps(response_obj).encode("utf-8")
    t_json_ms = (time.perf_counter() - t_json_start) * 1000.0

    payload_kb = len(json_bytes) / 1024.0
    t_total_ms = (time.perf_counter() - t_route_start) * 1000.0

    if ENABLE_PERFORMANCE_LOGS:
        print(f"[PERF-API] GET /batch-status: Server Proc={t_proc_ms:.2f}ms | JSON Serialization={t_json_ms:.2f}ms | Response Size={payload_kb:.2f} KB | Total={t_total_ms:.2f}ms")

    return Response(json_bytes, mimetype="application/json")


@app.route("/stop-batch/<session_id>", methods=["POST"])
def stop_batch(session_id):
    if session_id not in MEMORY_STORE:
        return jsonify({"error": "Invalid session ID"}), 404
    MEMORY_STORE[session_id]["stop"] = True
    return jsonify({"message": "Stop signal sent"})


@app.route("/batch-history", methods=["GET"])
def get_batch_history():
    return jsonify(BATCH_HISTORY)


@app.route("/update-question/<session_id>/<int:question_id>", methods=["PUT"])
def update_question(session_id, question_id):
    if session_id not in MEMORY_STORE:
        return jsonify({"error": "Invalid session ID"}), 404

    store = MEMORY_STORE[session_id]
    data = request.json
    new_question_text = data.get("question")

    if not new_question_text:
        return jsonify({"error": "Question text is required"}), 400

    # Find the question
    question_to_update = None
    for res in store["results"]:
        if res.get("id") == question_id:
            question_to_update = res
            break

    if not question_to_update:
        return jsonify({"error": "Question not found"}), 404

    # Re-classify
    cleaned_new = clean_source_question(new_question_text)
    question_to_update["question"] = cleaned_new
    question_to_update["original_question"] = new_question_text
    pred_bloom, diff, conf = classify_text(cleaned_new)
    question_to_update["predicted_bloom"] = pred_bloom
    question_to_update["difficulty"] = diff
    question_to_update["confidence"] = conf

    explanation = generate_dynamic_explanation(cleaned_new, "", pred_bloom, diff)
    question_to_update["explanation"] = explanation
    question_to_update["bloom_level"] = pred_bloom

    update_batch_history_stats(session_id)

    return jsonify(
        {
            "message": "Question updated successfully",
            "result": question_to_update,
        }
    )


@app.route("/export-rejections", methods=["GET"])
def export_rejections():
    if REJECTION_LOGS:
        with open("rejection_statistics.csv", "w", newline="", encoding="utf-8") as f:
            writer = csv.DictWriter(
                f, fieldnames=["question", "target_bloom", "reason"]
            )
            writer.writeheader()
            writer.writerows(REJECTION_LOGS)

    breakdown = {}
    for log in REJECTION_LOGS:
        r = log["reason"]
        breakdown[r] = breakdown.get(r, 0) + 1

    with open("rejection_breakdown.json", "w", encoding="utf-8") as f:
        json.dump(breakdown, f, indent=4)

    return jsonify(
        {
            "message": "Exported successfully",
            "csv": "rejection_statistics.csv",
            "json": "rejection_breakdown.json",
            "breakdown": breakdown,
        }
    )


@app.route("/delete-question/<session_id>/<int:question_id>", methods=["DELETE"])
def delete_question(session_id, question_id):
    if session_id not in MEMORY_STORE:
        return jsonify({"error": "Invalid session ID"}), 404

    store = MEMORY_STORE[session_id]

    initial_count = len(store["results"])
    store["results"] = [r for r in store["results"] if r.get("id") != question_id]

    if len(store["results"]) == initial_count:
        return jsonify({"error": "Question not found"}), 404

    store["total_questions"] = len(store["results"])
    store["completed_questions"] = len(store["results"])

    # Update BATCH_HISTORY if present
    for hist in BATCH_HISTORY:
        if hist["session_id"] == session_id:
            hist["total_questions"] = store["total_questions"]
            hist["completed_questions"] = store["completed_questions"]
            break

    update_batch_history_stats(session_id)

    return jsonify({"message": "Question deleted successfully"})





RE_CLEAN_PUNCT = re.compile(r'[^a-zA-Z0-9\s-]')

_STARTERS = tuple(sorted([
    "differentiate the performance characteristics of",
    "discuss the trade-offs between",
    "evaluate the suitability of",
    "evaluate the effectiveness of",
    "evaluate the trade-offs between",
    "assess the effectiveness of",
    "justify the selection of",
    "analyze the impact of",
    "analyze how",
    "analyze the",
    "explain the purpose of a",
    "explain the purpose of an",
    "explain the purpose of the",
    "explain the purpose of",
    "explain how",
    "explain why",
    "explain the concept of",
    "explain the process of",
    "explain",
    "describe the process of",
    "describe the",
    "describe",
    "what is the process of",
    "what is the concept of",
    "what is a",
    "what is an",
    "what is the",
    "what are the",
    "what is",
    "what are",
    "define",
    "apply the concept of",
    "apply",
    "evaluate",
    "design a secure",
    "design a",
    "design an",
    "design",
    "create a",
    "create an",
    "create",
    "write a",
    "implement",
    "determine how",
    "examine the",
    "examine",
    "interpret the",
    "interpret",
    "investigate the causes of",
    "investigate how",
    "investigate",
    "recommend a",
    "recommend an",
    "recommend",
    "validate whether",
    "formulate a",
    "critique the effectiveness of",
    "critique the",
    "critique",
    "optimize the",
    "optimize",
    "develop a",
    "develop",
    "propose a",
    "propose",
    "construct a",
    "construct",
    "differentiate",
    "assess",
    "justify",
    "how does a",
    "how does",
    "how do",
    "how",
    "discuss",
], key=len, reverse=True))

_PREFIXES = ("a ", "an ", "the ", "our ", "your ", "some ", "of ", "for ", "to ", "in ", "using ", "using a ", "using an ", "using the ")

_REMOVE_PHRASES = [
    "in enterprise databases", "enterprise databases", "database systems",
    "information systems", "software systems", "network systems",
    "distributed systems", "large-scale systems", "organizational environments",
    "for enterprise databases", "for database systems", "for organizations",
    "for enterprises", "in organizations", "in enterprises",
]
RE_REMOVE_PHRASES = re.compile('|'.join(re.escape(p) for p in sorted(_REMOVE_PHRASES, key=len, reverse=True)))

_GENERIC_PREFIXES = (
    "database solution", "software application", "network system",
    "system solution", "solution", "system", "framework",
    "application", "process",
)

_CONNECTORS = (
    "capable of handling", "capable of", "handling", "designed to",
    "designed for", "used for", "responsible for", "supports",
    "provides", "containing", "involving", "based on", "applying",
    "demonstrating", "consisting of", "relating to", "to", "for", "of",
)

_SPLIT_CLAUSES = [
    " and its role", " and how", " and why", " and the", " role in",
    " to reduce", " reducing", " in reducing", " versus", " vs",
    " improve", " improves", " improving", " perform", " performs",
    " performing", " use ", " uses ", " using ", " build ",
    " builds ", " building ", " work ", " works ", " working ",
]

_TRAILING_STOP_WORDS = {
    "and", "or", "how", "how do", "what", "what is", "the", "of", "for", "to", "a", "an", "in", "its", "role", "reducing", "redundancy"
}


@functools.lru_cache(maxsize=1024)
def extract_core_phrase(question: str) -> str:
    q_lower = question.lower().strip()

    # 1. Clean punctuation keeping hyphens
    q_clean = " ".join(RE_CLEAN_PUNCT.sub(" ", q_lower).split())

    # 2. Match and strip question/verb starters
    concept = q_clean
    for s in _STARTERS:
        if concept.startswith(s + " ") or concept == s:
            concept = concept[len(s):].strip()
            break

    # Strip leading articles and prepositions
    for prefix in _PREFIXES:
        if concept.startswith(prefix):
            concept = concept[len(prefix):].strip()
            break

    # Remove generic context words in single pass
    concept = RE_REMOVE_PHRASES.sub("", concept)
    concept = " ".join(concept.split()).strip()

    # 3. Detect generic prefixes followed by connectors
    matched_generic = False
    for gp in _GENERIC_PREFIXES:
        if concept.startswith(gp + " ") or concept == gp:
            rest = concept[len(gp):].strip()
            for conn in _CONNECTORS:
                if rest.startswith(conn + " ") or rest == conn:
                    candidate = rest[len(conn):].strip()
                    if candidate:
                        concept = candidate
                        matched_generic = True
                        break
            if matched_generic:
                break

    # 4. Remove split clauses and trailing stop words
    for clause in _SPLIT_CLAUSES:
        idx = concept.find(clause)
        if idx != -1:
            concept = concept[:idx]

    concept = " ".join(concept.split()).strip()

    words = concept.split()
    while words and words[-1] in _TRAILING_STOP_WORDS:
        words.pop()
    concept = " ".join(words).strip()

    if concept == "petabyte-scale transactions":
        concept = "petabyte-scale transaction processing"

    words = concept.split()
    if len(words) > 5:
        concept = " ".join(words[:5])

    return concept


def infer_domain(question: str, concept: str) -> str:
    """Infer the academic domain from the question or concept using the configured DOMAIN_MAP."""
    q_lower = question.lower()
    c_lower = concept.lower()

    for domain, keywords in DOMAIN_MAP.items():
        if domain in q_lower:
            return domain.title() if domain != "dbms" else "DBMS"

    for domain, keywords in DOMAIN_MAP.items():
        for kw in keywords:
            if kw in c_lower or kw in q_lower:
                if domain == "dbms":
                    return "DBMS"
                return domain.title()

    return "General Computer Science"


from question_understanding import QuestionUnderstandingEngine
from question_profile import QuestionProfile
from pipeline_context import PipelineContext

def pre_cache_embeddings(texts: list, st_model):
    if not texts or st_model is None:
        return
    
    unique_non_cached = []
    with EMBEDDING_CACHE_LOCK:
        for t in texts:
            key = normalize_embedding_key(t)
            if key not in EMBEDDING_CACHE:
                unique_non_cached.append((t, key))
                
    if not unique_non_cached:
        return
        
    try:
        raw_texts = [item[0] for item in unique_non_cached]
        embeddings = st_model.encode(raw_texts, convert_to_tensor=True)
        
        with EMBEDDING_CACHE_LOCK:
            for (t, key), emb in zip(unique_non_cached, embeddings):
                EMBEDDING_CACHE[key] = emb
    except Exception as e:
        print(f"Error pre-caching embeddings: {e}")


def classify_texts_batch(texts: list) -> list:
    if not texts:
        return []
    if deberta_model is None or deberta_tokenizer is None:
        raise RuntimeError("DeBERTa model is not loaded.")
    try:
        device = next(deberta_model.parameters()).device
        inputs = deberta_tokenizer(
            texts, padding=True, truncation=True, max_length=512, return_tensors="pt"
        )
        inputs = {k: v.to(device) for k, v in inputs.items()}
        with torch.inference_mode():
            outputs = deberta_model(**inputs)
            logits = outputs.logits
            probs = torch.nn.functional.softmax(logits, dim=-1)
            
            results = []
            index_to_bloom = {
                0: "Remember",
                1: "Understand",
                2: "Apply",
                3: "Analyze",
                4: "Evaluate",
                5: "Create",
            }
            for i in range(len(texts)):
                pred_class = logits[i].argmax().item()
                conf = probs[i][pred_class].item() * 100
                bloom = index_to_bloom.get(pred_class, "Analyze")
                diff = BLOOM_TO_DIFFICULTY.get(bloom, "Medium")
                results.append((bloom, diff, round(conf, 2)))
            return results
    except Exception as e:
        print(f"Error in batch classification: {e}")
        return [("Unknown", "Unknown", 0.0)] * len(texts)


def _write_benchmark_logs(all_candidates_log: list):
    try:
        log_file_path = "candidate_ranking_log.json"
        existing_candidates = []
        if os.path.exists(log_file_path):
            try:
                with open(log_file_path, "r", encoding="utf-8") as lf:
                    existing_candidates = json.load(lf)
            except Exception:
                existing_candidates = []
                
        for c in all_candidates_log:
            existing_candidates.append({
                "Candidate ID": c.get("Candidate ID", ""),
                "Question": c.get("Question", ""),
                "Bloom Prediction": c.get("Bloom Prediction", ""),
                "Difficulty": c.get("Difficulty", ""),
                "Confidence": c.get("Confidence", 0.0),
                "Concept Score": c.get("Concept Score", ""),
                "Duplicate Score": c.get("Duplicate Score", 0.0),
                "Entity Score": c.get("Entity Score", 0.0),
                "Validation Status": c.get("Validation Status", ""),
                "Rejection Reason": c.get("Rejection Reason", ""),
                "Generation Round": c.get("Generation Round", 1),
                "Rank Score": c.get("Rank Score", 0.0),
            })

        with open(log_file_path, "w", encoding="utf-8") as lf:
            json.dump(existing_candidates, lf, indent=4)
    except Exception as exc:
        print(f"[WARNING] Error writing to candidate_ranking_log.json: {exc}")


def generate_validated_variant(
    question,
    src_bloom,
    src_diff,
    target_bloom,
    target_difficulty,
    domain,
    required_concept,
    session_seen,
    config_mode=None,
    mode_name=None,
):
    """
    Executes the unified Mode E generation pipeline with batched model evaluation.
    Applies all validation checks and returns the best candidate or the best failed attempt.
    """
    start_time = time.time()
    t_pipeline_start = time.perf_counter()
    st_hits_start = EMBEDDING_CACHE_HITS
    st_misses_start = EMBEDDING_CACHE_MISSES

    t_tok = flan_tokenizer
    
    if config_mode is None:
        config_mode = MODES[ACTIVE_MODE]
        mode_name = ACTIVE_MODE
    if mode_name is None:
        mode_name = ACTIVE_MODE

    # Timing accumulators for performance profiling
    cum_qu = 0.0
    cum_dom = 0.0
    cum_conc = 0.0
    cum_prompt = 0.0
    cum_flan = 0.0
    cum_dedup = 0.0
    cum_deberta = 0.0
    cum_prof = 0.0
    cum_st = 0.0
    cum_val = 0.0
    cum_rank = 0.0
    cum_exp = 0.0
    cum_retry = 0.0
    st_call_count = 0
        
    t_qu_start = time.perf_counter()
    orig_profile = QuestionUnderstandingEngine.build_profile(question, src_bloom)
    t_qu_ms = (time.perf_counter() - t_qu_start) * 1000.0
    cum_qu += t_qu_ms
    prof_duration = t_qu_ms / 1000.0
    
    pipeline_ctx = PipelineContext(
        source_question=question,
        source_bloom=src_bloom,
        target_bloom=target_bloom,
        target_difficulty=target_difficulty,
        source_profile=orig_profile,
    )
    pipeline_ctx.timings["profiling"] = prof_duration
    
    # Extract domain and concept from the profile
    t_dom_start = time.perf_counter()
    domain = orig_profile.domain
    t_dom_ms = (time.perf_counter() - t_dom_start) * 1000.0
    cum_dom += t_dom_ms

    t_conc_start = time.perf_counter()
    required_concept = orig_profile.topic
    t_conc_ms = (time.perf_counter() - t_conc_start) * 1000.0
    cum_conc += t_conc_ms
    
    # Pre-cache original question/concepts in SentenceTransformer embedding cache
    st_model = _get_st_model()
    t_st_init_ms = 0.0
    if st_model is not None:
        t_st_s = time.perf_counter()
        orig_texts = [orig_profile.normalized_question, orig_profile.raw_question]
        orig_texts.extend(orig_profile.concepts)
        orig_texts.extend(orig_profile.technical_entities)
        orig_texts.extend(orig_profile.noun_chunks)
        pre_cache_embeddings(orig_texts, st_model)
        t_st_init_ms = (time.perf_counter() - t_st_s) * 1000.0
        cum_st += t_st_init_ms
        st_call_count += 1
        
    flan_calls_count = 0
    deberta_calls_count = 0
    candidates_generated_count = 0
    candidates_validated_count = 0
    candidates_rejected_before_validation_count = 0
    
    all_candidates_log = []
    seen_questions = set()
    best_candidate = None
    retry_history_log = []
    prev_best_score = -1.0
    is_semantic_drift_round = False
    failure_hint_round_2 = ""
    rejected_questions = []
    
    max_rounds = config_mode.get("max_generation_rounds", 2)
    
    for current_round in range(1, max_rounds + 1):
        t_round_start = time.perf_counter()
        flan_calls_count += 1
        
        candidates_list, prompt_used, t_prompt_ms, t_flan_ms = generate_single_variant(
            question,
            src_bloom,
            src_diff,
            target_bloom,
            target_difficulty,
            domain=domain,
            topic=required_concept,
            required_concept=required_concept,
            attempt_number=current_round,
            failure_hint=failure_hint_round_2 if current_round > 1 else "",
            config_mode=config_mode,
            mode_name=mode_name,
            return_timings=True,
        )
        cum_prompt += t_prompt_ms
        cum_flan += t_flan_ms
        
        if not isinstance(candidates_list, list):
            candidates_list = [candidates_list]
            
        candidates_generated_count += len(candidates_list)
        
        # Deduplicate
        t_dedup_start = time.perf_counter()
        unique_candidates = []
        for idx, cand_text in enumerate(candidates_list):
            cand_text = sanitize_generated_question(cand_text)
            
            if cand_text in unique_candidates or cand_text in session_seen or cand_text in seen_questions:
                candidates_rejected_before_validation_count += 1
                candidate_info = {
                    "Candidate ID": f"Round_{current_round}_Index_{idx}",
                    "Question": cand_text,
                    "Bloom Prediction": "N/A",
                    "Difficulty": "N/A",
                    "Confidence": 0.0,
                    "Concept Score": "N/A",
                    "Duplicate Score": 1.0,
                    "Entity Score": 0.0,
                    "entity_score": 0.0,
                    "Validation Status": "Deduplicated",
                    "Rejection Reason": "Duplicate",
                    "Generation Round": current_round,
                    "Rank Score": 0.0,
                    
                    "bloom_prediction": "N/A",
                    "difficulty_prediction": "N/A",
                    "confidence": 0.0,
                    "rejection_reason": "Duplicate",
                    "concept_similarity_score": 0.0,
                    "duplicate_score": 1.0,
                    "validation_status": "Deduplicated",
                    "generated_question": cand_text,
                    "attempt_number": current_round,
                    "prompt": prompt_used,
                    "prompt_tokens": len(t_tok.encode(prompt_used)) if t_tok is not None else 0,
                    "completion_tokens": len(t_tok.encode(cand_text)) if t_tok is not None else 0,
                    "concept_match_score": "N/A",
                    "concept_match_method": "N/A",
                    "concept_similarity_score_raw": 0.0,
                    "matched_keywords": "",
                    
                    "bloom_score": 0.0,
                    "concept_score": 0.0,
                    "entity_score": 0.0,
                    "semantic_score": 0.0,
                    "number_score": 0.0,
                    "grammar_score": 0.0,
                    "duplicate_score": 0.0,
                    "domain_score": 0.0,
                    "topic_score": 0.0,
                    "total_score": 0.0,
                }
                all_candidates_log.append(candidate_info)
            else:
                unique_candidates.append(cand_text)
        t_dedup_ms = (time.perf_counter() - t_dedup_start) * 1000.0
        cum_dedup += t_dedup_ms
                
        if not unique_candidates:
            t_round_total_ms = (time.perf_counter() - t_round_start) * 1000.0
            if current_round > 1:
                cum_retry += t_round_total_ms
            if ENABLE_GENERATION_PERFORMANCE_LOGS:
                print("==================================================", flush=True)
                print(f"GENERATION ROUND {current_round}", flush=True)
                print("==================================================", flush=True)
                print(f"Question Understanding : {t_qu_ms if current_round==1 else 0.0:.2f} ms", flush=True)
                print(f"Domain Detection       : {t_dom_ms if current_round==1 else 0.0:.2f} ms", flush=True)
                print(f"Concept Extraction     : {t_conc_ms if current_round==1 else 0.0:.2f} ms", flush=True)
                print(f"Prompt Construction    : {t_prompt_ms:.2f} ms", flush=True)
                print(f"FLAN Generation        : {t_flan_ms:.2f} ms", flush=True)
                print(f"Deduplication          : {t_dedup_ms:.2f} ms", flush=True)
                print("Batch DeBERTa          : 0.00 ms", flush=True)
                print("Profile Creation       : 0.00 ms", flush=True)
                print("Sentence Embeddings    : 0.00 ms", flush=True)
                print("Validation             : 0.00 ms", flush=True)
                print("Ranking                : 0.00 ms", flush=True)
                print("Explanation            : 0.00 ms", flush=True)
                print(f"\nRound Total            : {t_round_total_ms:.2f} ms", flush=True)
                print("==================================================\n", flush=True)
            continue
            
        # Batch pre-computations
        t_deb_start = time.perf_counter()
        batch_classifications = classify_texts_batch(unique_candidates)
        t_deb_ms = (time.perf_counter() - t_deb_start) * 1000.0
        cum_deberta += t_deb_ms

        classification_by_question = dict(zip(unique_candidates, batch_classifications))
        
        def cached_classify_deberta(text: str):
            if text in classification_by_question:
                return classification_by_question[text]
            return classify_text(text)
            
        t_prof_start = time.perf_counter()
        cand_profiles_temp = []
        texts_to_pre_cache = []
        for gen_text in unique_candidates:
            cp = QuestionUnderstandingEngine.build_profile(gen_text, target_bloom)
            cand_profiles_temp.append((gen_text, cp))
            texts_to_pre_cache.append(cp.normalized_question)
            texts_to_pre_cache.extend(cp.concepts)
            texts_to_pre_cache.extend(cp.technical_entities)
            texts_to_pre_cache.extend(cp.noun_chunks)
        t_prof_ms = (time.perf_counter() - t_prof_start) * 1000.0
        cum_prof += t_prof_ms
            
        t_st_r_start = time.perf_counter()
        if st_model is not None:
            pre_cache_embeddings(texts_to_pre_cache, st_model)
            st_call_count += 1
        t_st_r_ms = (time.perf_counter() - t_st_r_start) * 1000.0
        cum_st += t_st_r_ms
            
        # Validate unique candidates
        t_val_start = time.perf_counter()
        round_candidates = []
        is_final_retries = (current_round == max_rounds)
        
        for idx, (gen_text, cand_prof) in enumerate(cand_profiles_temp):
            pipeline_ctx.generated_candidates.append(gen_text)
            pipeline_ctx.candidate_profiles.append(cand_prof)
            
            val_out = evaluate_candidate(
                original_q=orig_profile,
                candidate_q=cand_prof,
                target_bloom=target_bloom,
                target_difficulty=target_difficulty,
                seen_questions=seen_questions,
                session_seen=session_seen,
                config_mode=config_mode,
                deberta_classifier_fn=cached_classify_deberta,
                st_model=st_model,
                get_cached_embedding_fn=get_cached_embedding,
                validate_bloom_verbs_fn=validate_bloom_verbs,
                is_final_round=is_final_retries
            )
            pipeline_ctx.validation_results.append(val_out)
            
            new_bloom = val_out.detailed_metrics.get("bloom", {}).get("predicted_bloom", "Unknown") if "bloom" in val_out.detailed_metrics else "N/A"
            new_diff = val_out.detailed_metrics.get("bloom", {}).get("predicted_difficulty", "Unknown") if "bloom" in val_out.detailed_metrics else "N/A"
            new_conf = val_out.detailed_metrics.get("bloom", {}).get("confidence", 0.0) if "bloom" in val_out.detailed_metrics else 0.0
            
            concept_match_score = f"{len(val_out.detailed_metrics.get('concept', {}).get('matched_concepts', []))}/{len(val_out.detailed_metrics.get('concept', {}).get('original_concepts', []))}" if "concept" in val_out.detailed_metrics else "N/A"
            concept_sim = val_out.detailed_metrics.get("concept", {}).get("preservation_percentage", 0.0) if "concept" in val_out.detailed_metrics else "N/A"
            concept_method = "nlp_spacy"
            matched_words_str = ", ".join(val_out.detailed_metrics.get("concept", {}).get("matched_concepts", [])) if "concept" in val_out.detailed_metrics else ""
            
            max_duplicate_ratio = val_out.detailed_metrics.get("duplicate", {}).get("max_duplicate_ratio", 0.0) if "duplicate" in val_out.detailed_metrics else 1.0
            
            if val_out.rejection_reason not in ("Too Short", "Duplicate"):
                candidates_validated_count += 1
                if "bloom" in val_out.detailed_metrics:
                    deberta_calls_count += 1
                    
            candidate_info = {
                "Candidate ID": f"Round_{current_round}_Index_{idx}",
                "Question": gen_text,
                "Bloom Prediction": new_bloom,
                "Difficulty": new_diff,
                "Confidence": round(new_conf, 2),
                "Concept Score": concept_match_score,
                "Duplicate Score": round(max_duplicate_ratio, 4),
                "Entity Score": val_out.entity_score,
                "entity_score": val_out.entity_score,
                "Validation Status": "Pass" if val_out.passed else "Fail",
                "Rejection Reason": val_out.rejection_reason,
                "Generation Round": current_round,
                "Rank Score": calculate_nlp_rank_score(val_out),
                
                # Compatibility fields
                "bloom_prediction": new_bloom,
                "difficulty_prediction": new_diff,
                "confidence": new_conf,
                "rejection_reason": val_out.rejection_reason,
                "concept_similarity_score": concept_sim,
                "duplicate_score": max_duplicate_ratio,
                "validation_status": "Pass" if val_out.passed else "Fail",
                "generated_question": gen_text,
                "attempt_number": current_round,
                "prompt": prompt_used,
                "prompt_tokens": len(t_tok.encode(prompt_used)) if t_tok is not None else 0,
                "completion_tokens": len(t_tok.encode(gen_text)) if t_tok is not None else 0,
                "concept_match_score": concept_match_score,
                "concept_match_method": concept_method,
                "concept_similarity_score_raw": val_out.detailed_metrics.get("semantic", {}).get("similarity", 0.0) if "semantic" in val_out.detailed_metrics else 0.0,
                "matched_keywords": matched_words_str,
                
                # Per-stage scores for benchmark reporting
                "bloom_score": val_out.bloom_score,
                "concept_score": val_out.concept_score,
                "entity_score": val_out.entity_score,
                "semantic_score": val_out.semantic_score,
                "number_score": val_out.number_score,
                "grammar_score": val_out.grammar_score,
                "duplicate_score": val_out.duplicate_score,
                "domain_score": val_out.domain_score,
                "topic_score": val_out.topic_score,
                "total_score": val_out.total_score,
                "val_out": val_out,
                "cand_prof": cand_prof,
            }
            
            round_candidates.append(candidate_info)
            all_candidates_log.append(candidate_info)
            if not val_out.passed:
                rejected_questions.append(gen_text)
        t_val_ms = (time.perf_counter() - t_val_start) * 1000.0
        cum_val += t_val_ms
        
        t_rank_start = time.perf_counter()
        passing_in_round = [c for c in round_candidates if c["validation_status"] == "Pass"]
        if passing_in_round:
            passing_in_round = rank_candidates_dicts(passing_in_round)
            best_candidate = passing_in_round[0]
            
            best_round_score = best_candidate.get("total_score", 100.0)
            improved = False
            if current_round > 1:
                improved = best_round_score > prev_best_score
            retry_history_log.append({
                "round": current_round,
                "failure_reason": "None",
                "retry_feedback": "",
                "validation_score": best_round_score,
                "passed": True,
                "improved_from_previous": improved
            })
            t_rank_ms = (time.perf_counter() - t_rank_start) * 1000.0
            cum_rank += t_rank_ms

            t_round_total_ms = (time.perf_counter() - t_round_start) * 1000.0
            if current_round > 1:
                cum_retry += t_round_total_ms

            if ENABLE_GENERATION_PERFORMANCE_LOGS:
                print("==================================================", flush=True)
                print(f"GENERATION ROUND {current_round}", flush=True)
                print("==================================================", flush=True)
                print(f"Question Understanding : {t_qu_ms if current_round==1 else 0.0:.2f} ms", flush=True)
                print(f"Domain Detection       : {t_dom_ms if current_round==1 else 0.0:.2f} ms", flush=True)
                print(f"Concept Extraction     : {t_conc_ms if current_round==1 else 0.0:.2f} ms", flush=True)
                print(f"Prompt Construction    : {t_prompt_ms:.2f} ms", flush=True)
                print(f"FLAN Generation        : {t_flan_ms:.2f} ms", flush=True)
                print(f"Deduplication          : {t_dedup_ms:.2f} ms", flush=True)
                print(f"Batch DeBERTa          : {t_deb_ms:.2f} ms", flush=True)
                print(f"Profile Creation       : {t_prof_ms:.2f} ms", flush=True)
                print(f"Sentence Embeddings    : {t_st_r_ms:.2f} ms", flush=True)
                print(f"Validation             : {t_val_ms:.2f} ms", flush=True)
                print(f"Ranking                : {t_rank_ms:.2f} ms", flush=True)
                print("Explanation            : 0.00 ms", flush=True)
                print(f"\nRound Total            : {t_round_total_ms:.2f} ms", flush=True)
                print("==================================================\n", flush=True)
            break
        else:
            for c in round_candidates:
                seen_questions.add(c["generated_question"])
                
            round_rejection_reasons = set()
            best_round_score = -1.0
            
            for c in round_candidates:
                vo = c.get("val_out")
                if vo:
                    round_rejection_reasons.add(vo.rejection_reason)
                    best_round_score = max(best_round_score, vo.total_score)
                    
            improved = False
            if current_round > 1:
                improved = best_round_score > prev_best_score
            prev_best_score = best_round_score
            
            has_drift = "Semantic Drift" in round_rejection_reasons or "Concept Drift" in round_rejection_reasons
            has_missing_or_reduced = False
            for c in round_candidates:
                vo = c.get("val_out")
                if vo and hasattr(vo, "detailed_metrics") and vo.detailed_metrics:
                    concept_data = vo.detailed_metrics.get("concept", {}) or {}
                    missing_concepts = concept_data.get("missing_concepts", []) or []
                    concept_preservation = concept_data.get("preservation_percentage", 1.0)
                    
                    entity_data = vo.detailed_metrics.get("entity", {}) or {}
                    missing_entities = entity_data.get("missing_entities", []) or []
                    
                    if len(missing_concepts) > 0 or len(missing_entities) > 0 or concept_preservation < 1.0:
                        has_missing_or_reduced = True
                        break
            is_semantic_drift_round = has_drift and has_missing_or_reduced
            
            failure_reasons_str = ", ".join(sorted(list(round_rejection_reasons))) if round_rejection_reasons else "None"
            retry_history_log.append({
                "round": current_round,
                "failure_reason": failure_reasons_str,
                "retry_feedback": "",
                "validation_score": best_round_score,
                "passed": False,
                "improved_from_previous": improved
            })
            t_rank_ms = (time.perf_counter() - t_rank_start) * 1000.0
            cum_rank += t_rank_ms

            t_round_total_ms = (time.perf_counter() - t_round_start) * 1000.0
            if current_round > 1:
                cum_retry += t_round_total_ms

            if ENABLE_GENERATION_PERFORMANCE_LOGS:
                print("==================================================", flush=True)
                print(f"GENERATION ROUND {current_round}", flush=True)
                print("==================================================", flush=True)
                print(f"Question Understanding : {t_qu_ms if current_round==1 else 0.0:.2f} ms", flush=True)
                print(f"Domain Detection       : {t_dom_ms if current_round==1 else 0.0:.2f} ms", flush=True)
                print(f"Concept Extraction     : {t_conc_ms if current_round==1 else 0.0:.2f} ms", flush=True)
                print(f"Prompt Construction    : {t_prompt_ms:.2f} ms", flush=True)
                print(f"FLAN Generation        : {t_flan_ms:.2f} ms", flush=True)
                print(f"Deduplication          : {t_dedup_ms:.2f} ms", flush=True)
                print(f"Batch DeBERTa          : {t_deb_ms:.2f} ms", flush=True)
                print(f"Profile Creation       : {t_prof_ms:.2f} ms", flush=True)
                print(f"Sentence Embeddings    : {t_st_r_ms:.2f} ms", flush=True)
                print(f"Validation             : {t_val_ms:.2f} ms", flush=True)
                print(f"Ranking                : {t_rank_ms:.2f} ms", flush=True)
                print("Explanation            : 0.00 ms", flush=True)
                print(f"\nRound Total            : {t_round_total_ms:.2f} ms", flush=True)
                print("==================================================\n", flush=True)

        if current_round < max_rounds:
            recent_rejected = rejected_questions[-5:]
            rejected_questions_list = [f"{idx_rej+1}. {q_text}" for idx_rej, q_text in enumerate(recent_rejected)]
            rejected_str = "\n".join(rejected_questions_list)
            
            verbs_map = {
                "Create": "Design, Create, Develop, Formulate, or Construct",
                "Evaluate": "Evaluate, Assess, Critique, Justify, Defend, or Determine",
                "Analyze": "Analyze, Compare, Differentiate, Examine, Contrast, or Distinguish",
                "Apply": "Apply, Implement, Use, Demonstrate, Solve, or Calculate",
                "Understand": "Explain, Describe, Discuss, Summarize, or Interpret",
                "Remember": "Define, Identify, List, Name, State, or Recall"
            }
            allowed_verbs_str = verbs_map.get(target_bloom, "")
            
            failure_hint_round_2 = (
                f"\n\nPrevious rejected questions:\n{rejected_str}\n\n"
                f"Avoid the structures and verbs of the rejected questions above. "
                f"The new question MUST start with one of the following verbs: {allowed_verbs_str}. "
                f"Ensure it evaluates the concept '{required_concept}' at the {target_bloom} ({target_difficulty}) level."
            )
            
    generation_time = time.time() - start_time
    
    # Decoupled Logging I/O
    _write_benchmark_logs(all_candidates_log)
    
    # Build attempts list schema
    attempts_list = []
    for c in all_candidates_log:
        if c.get("validation_status") in ["Pass", "Fail"]:
            attempts_list.append({
                "attempt": c.get("attempt_number", 1),
                "attempt_number": c.get("attempt_number", 1),
                "generated_question": c.get("Question", ""),
                "question": c.get("Question", ""),
                "bloom_prediction": c.get("bloom_prediction", ""),
                "predicted_bloom": c.get("bloom_prediction", ""),
                "difficulty_prediction": c.get("difficulty_prediction", ""),
                "predicted_difficulty": c.get("difficulty_prediction", ""),
                "confidence": c.get("confidence", 0.0),
                "concept_score": c.get("concept_match_score", 0.0),
                "duplicate_score": c.get("duplicate_score", 0.0),
                "entity_score": c.get("entity_score", 0.0),
                "bloom_score": c.get("bloom_score", 0.0),
                "semantic_score": c.get("semantic_score", 0.0),
                "number_score": c.get("number_score", 0.0),
                "grammar_score": c.get("grammar_score", 0.0),
                "domain_score": c.get("domain_score", 0.0),
                "topic_score": c.get("topic_score", 0.0),
                "total_score": c.get("total_score", 0.0),
                "validation_status": c.get("validation_status", ""),
                "rejection_reason": c.get("rejection_reason", ""),
                "generation_time": generation_time / flan_calls_count if flan_calls_count else 0.0,
                "val_out": c.get("val_out"),
                "cand_prof": c.get("cand_prof"),
                "prompt": c.get("prompt", ""),
                "prompt_tokens": c.get("prompt_tokens", 0),
                "completion_tokens": c.get("completion_tokens", 0),
                "advisor_activated": False,
                "advisor_guidance_produced": False,
                "advisor_skipped": False,
                "missing_preferred": [],
                "missing_supporting": [],
                "generic_terms_detected": [],
            })
            
    if best_candidate is None:
        t_rank_final_start = time.perf_counter()
        all_validated = [c for c in all_candidates_log if c.get("validation_status") in ["Pass", "Fail"]]
        if all_validated:
            all_validated = rank_candidates_dicts(all_validated)
            best_candidate = all_validated[0]
        t_rank_final_ms = (time.perf_counter() - t_rank_final_start) * 1000.0
        cum_rank += t_rank_final_ms
            
    if best_candidate is None:
        return ValidationResult(
            generated_question="[FAILED]",
            source_question=question,
            source_bloom=src_bloom,
            source_difficulty=src_diff,
            target_bloom=target_bloom,
            target_difficulty=target_difficulty,
            predicted_bloom="Unknown",
            predicted_difficulty="Unknown",
            confidence=0.0,
            attempts=flan_calls_count,
            generation_time=generation_time,
            concept_match_score="N/A",
            rejection_reason="Fatal Error",
            prompt_used="",
            explanation="Fatal error: No generation attempt occurred.",
            validation_status="Fail",
            attempts_list=[],
            flan_calls=flan_calls_count,
            deberta_calls=deberta_calls_count,
            candidates_generated=candidates_generated_count,
            candidates_validated=candidates_validated_count,
            candidates_rejected_before_validation=candidates_rejected_before_validation_count,
            retry_history=retry_history_log,
        )
        
    is_pass = best_candidate.get("validation_status") == "Pass"
    gen_q = best_candidate.get("Question", "[FAILED]")
    p_bloom = best_candidate.get("bloom_prediction", "Unknown")
    p_diff = best_candidate.get("difficulty_prediction", "Unknown")
    
    t_exp_start = time.perf_counter()
    explanation = ""
    val_status = "Fail"
    
    if is_pass:
        explanation = generate_dynamic_explanation(gen_q, required_concept, p_bloom, p_diff)
        val_status = calculate_validation_status(
            p_bloom, p_diff, target_bloom, target_difficulty
        )
        session_seen.append(gen_q)
        max_history = config_mode["session_history_length"]
        while len(session_seen) > max_history:
            session_seen.pop(0)
    else:
        val_status = "Best Candidate"
        explanation_lines = [
            "<div class='text-amber-600 dark:text-amber-400 font-bold mb-2'>[VALIDATION FAILED - RETURNING BEST CANDIDATE]</div>",
            f"Generation failed after {flan_calls_count} rounds. Below is the diagnostic report:<br>"
        ]
        for att in attempts_list:
            explanation_lines.append(
                "<div class='mt-2.5 p-2 bg-slate-50 dark:bg-slate-800/50 "
                "border border-slate-100 dark:border-slate-800 rounded text-xs'>"
                f"<strong>Attempt {att['attempt_number']}</strong> &bull; "
                f"<span class='text-rose-500 font-semibold'>{att['rejection_reason']}</span><br>"
                f"<span class='text-slate-500'>Question:</span> \"{att['question']}\"<br>"
                f"<span class='text-slate-500'>Extracted Concept:</span> "
                f"\"{att.get('extracted_concept', required_concept)}\" &bull; "
                f"<span class='text-slate-500'>Concept Match:</span> "
                f"{att.get('concept_match_method', 'N/A')} "
                f"(similarity: {att.get('concept_similarity_score', 0.0)})<br>"
                f"<span class='text-slate-500'>Prediction:</span> "
                f"{att.get('predicted_bloom', 'Unknown')} "
                f"({att.get('predicted_difficulty', 'Unknown')}) &bull; "
                f"<span class='text-slate-500'>Confidence:</span> {att.get('confidence', 0.0)}% &bull; "
                f"<span class='text-slate-500'>Latency:</span> {att.get('generation_time', 0.0)}s"
                f"</div>"
            )
        explanation = "\n".join(explanation_lines)
    
    t_exp_ms = (time.perf_counter() - t_exp_start) * 1000.0
    cum_exp += t_exp_ms
        
    pipeline_ctx.best_candidate = gen_q
    t_pipeline_total_ms = (time.perf_counter() - t_pipeline_start) * 1000.0
    pipeline_ctx.timings["total_pipeline"] = t_pipeline_total_ms / 1000.0
    
    if ENABLE_GENERATION_PERFORMANCE_LOGS:
        stages_timing = {
            "Question Understanding": cum_qu,
            "Domain Detection": cum_dom,
            "Concept Extraction": cum_conc,
            "Prompt Construction": cum_prompt,
            "FLAN Generation": cum_flan,
            "Deduplication": cum_dedup,
            "Batch DeBERTa": cum_deberta,
            "Profile Creation": cum_prof,
            "Sentence Embeddings": cum_st,
            "Validation": cum_val,
            "Ranking": cum_rank,
            "Explanation": cum_exp,
        }
        slowest_stage = max(stages_timing, key=stages_timing.get)

        total_t = t_pipeline_total_ms if t_pipeline_total_ms > 0 else 1.0
        pct_qu = (cum_qu / total_t) * 100.0
        pct_dom = (cum_dom / total_t) * 100.0
        pct_conc = (cum_conc / total_t) * 100.0
        pct_prompt = (cum_prompt / total_t) * 100.0
        pct_flan = (cum_flan / total_t) * 100.0
        pct_deberta = (cum_deberta / total_t) * 100.0
        pct_val = (cum_val / total_t) * 100.0
        pct_rank = (cum_rank / total_t) * 100.0
        pct_exp = (cum_exp / total_t) * 100.0

        st_hits = EMBEDDING_CACHE_HITS - st_hits_start
        st_misses = EMBEDDING_CACHE_MISSES - st_misses_start

        passed_count = len([c for c in all_candidates_log if c.get("validation_status") == "Pass"])
        rejected_count = candidates_generated_count - passed_count
        retries_count = flan_calls_count - 1 if flan_calls_count > 0 else 0

        print("\n==================================================", flush=True)
        print("GENERATION PERFORMANCE SUMMARY", flush=True)
        print("==================================================", flush=True)
        print(f"Question Understanding : {cum_qu:.2f} ms", flush=True)
        print(f"Domain Detection       : {cum_dom:.2f} ms", flush=True)
        print(f"Concept Extraction     : {cum_conc:.2f} ms", flush=True)
        print(f"Prompt Construction    : {cum_prompt:.2f} ms", flush=True)
        print(f"FLAN Generation        : {cum_flan:.2f} ms", flush=True)
        print(f"Deduplication          : {cum_dedup:.2f} ms", flush=True)
        print(f"Batch DeBERTa          : {cum_deberta:.2f} ms", flush=True)
        print(f"Profile Creation       : {cum_prof:.2f} ms", flush=True)
        print(f"Sentence Embeddings    : {cum_st:.2f} ms", flush=True)
        print(f"Validation             : {cum_val:.2f} ms", flush=True)
        print(f"Ranking                : {cum_rank:.2f} ms", flush=True)
        print(f"Explanation            : {cum_exp:.2f} ms", flush=True)
        print(f"\nRetry Overhead         : {cum_retry:.2f} ms", flush=True)
        print(f"\nTOTAL PIPELINE         : {t_pipeline_total_ms:.2f} ms", flush=True)
        print(f"\nSlowest Stage          : {slowest_stage}", flush=True)
        print(f"\nPercentage of Total Time\n", flush=True)
        print(f"Question Understanding : {pct_qu:.2f}%", flush=True)
        print(f"Domain Detection       : {pct_dom:.2f}%", flush=True)
        print(f"Concept Extraction     : {pct_conc:.2f}%", flush=True)
        print(f"Prompt Construction    : {pct_prompt:.2f}%", flush=True)
        print(f"FLAN Generation        : {pct_flan:.2f}%", flush=True)
        print(f"Batch DeBERTa          : {pct_deberta:.2f}%", flush=True)
        print(f"Validation             : {pct_val:.2f}%", flush=True)
        print(f"Ranking                : {pct_rank:.2f}%", flush=True)
        print(f"Explanation            : {pct_exp:.2f}%\n", flush=True)
        print("==================================================", flush=True)
        print(f"• Number of FLAN generations               : {flan_calls_count}", flush=True)
        print(f"• Number of DeBERTa classifications         : {deberta_calls_count}", flush=True)
        print(f"• Number of generated candidates           : {candidates_generated_count}", flush=True)
        print(f"• Number of validated candidates           : {candidates_validated_count}", flush=True)
        print(f"• Number of rejected candidates            : {rejected_count}", flush=True)
        print(f"• Number of retries                        : {retries_count}", flush=True)
        print(f"• Number of SentenceTransformer calls       : {st_call_count}", flush=True)
        print(f"• Cache hits                               : {st_hits}", flush=True)
        print(f"• Cache misses                             : {st_misses}\n", flush=True)
    
    return ValidationResult(
        generated_question=gen_q,
        source_question=question,
        source_bloom=src_bloom,
        source_difficulty=src_diff,
        target_bloom=target_bloom,
        target_difficulty=target_difficulty,
        predicted_bloom=p_bloom,
        predicted_difficulty=p_diff,
        confidence=best_candidate.get("confidence", 0.0),
        attempts=flan_calls_count,
        generation_time=generation_time,
        concept_match_score=str(best_candidate.get("concept_match_score", 0.0)),
        rejection_reason=best_candidate.get("rejection_reason", "None"),
        prompt_used=best_candidate.get("prompt", ""),
        explanation=explanation,
        validation_status=val_status,
        attempts_list=attempts_list,
        prompt_tokens=best_candidate.get("prompt_tokens", 0),
        completion_tokens=best_candidate.get("completion_tokens", 0),
        flan_calls=flan_calls_count,
        deberta_calls=deberta_calls_count,
        candidates_generated=candidates_generated_count,
        candidates_validated=candidates_validated_count,
        candidates_rejected_before_validation=candidates_rejected_before_validation_count,
        retry_history=retry_history_log,
    )



@app.route("/rephrase", methods=["POST"])
def rephrase():
    data = request.json
    question = data.get("question")
    target_difficulty = data.get("target_difficulty")

    if not question or not target_difficulty:
        return (
            jsonify({"error": "Question and target difficulty required."}),
            400,
        )

    # Normalise legacy labels
    _legacy_map = {"Moderate": "Medium", "Difficult": "Hard"}
    target_difficulty = _legacy_map.get(target_difficulty, target_difficulty)

    if target_difficulty not in DIFFICULTY_TO_BLOOM:
        target_difficulty = "Medium"

    cleaned_q = clean_source_question(question)
    required_concept = normalize_academic_concept(cleaned_q)
    src_bloom, src_diff, _ = classify_text(cleaned_q)
    target_blooms = DIFFICULTY_TO_BLOOM[target_difficulty]

    if "seen_variants" not in session:
        session["seen_variants"] = []

    session_seen = session.get("seen_variants", [])
    domain = infer_domain(cleaned_q, required_concept)

    variants_to_return = []

    for target_bloom in target_blooms:
        # Call the centralized pipeline
        result = generate_validated_variant(
            question=cleaned_q,
            src_bloom=src_bloom,
            src_diff=src_diff,
            target_bloom=target_bloom,
            target_difficulty=target_difficulty,
            domain=domain,
            required_concept=required_concept,
            session_seen=session_seen,
        )

        # Update session
        session["seen_variants"] = session_seen
        session.modified = True

        if result.validation_status == "Fail":
            variants_to_return.append(
                {
                    "question": "[FAILED TO GENERATE VALID QUESTION]",
                    "target_bloom": target_bloom,
                    "target_difficulty": target_difficulty,
                    "predicted_bloom": "Unknown",
                    "predicted_difficulty": "Unknown",
                    "confidence": 0.0,
                    "validation_status": "Rejected",
                    "explanation": "No valid question satisfied Bloom level, difficulty, and confidence requirements.",
                }
            )
        else:
            variants_to_return.append(
                {
                    "question": result.generated_question,
                    "target_bloom": result.target_bloom,
                    "target_difficulty": result.target_difficulty,
                    "predicted_bloom": result.predicted_bloom,
                    "predicted_difficulty": result.predicted_difficulty,
                    "confidence": result.confidence,
                    "validation_status": result.validation_status,
                    "explanation": result.explanation,
                }
            )

    if not variants_to_return:
        return jsonify({"error": "Failed to generate variants."}), 500

    return jsonify({"original_question": question, "variants": variants_to_return})


@app.route("/export", methods=["POST"])
def export_data():
    format_type = request.form.get("format")
    session_id = request.form.get("session_id") or session.get("batch_id")

    if not session_id or session_id not in MEMORY_STORE:
        return (
            jsonify(
                {
                    "error": "No active session data found. Please classify a batch first."
                }
            ),
            404,
        )

    store = MEMORY_STORE[session_id]
    results = store.get("results", [])

    if not results:
        return jsonify({"error": "No results to export yet."}), 400

    if format_type == "csv":
        df = pd.DataFrame(results)
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], f"export_{session_id}.csv")
        df.to_csv(filepath, index=False)
        return send_file(
            filepath, as_attachment=True, download_name="bloom_results.csv"
        )
    elif format_type == "excel":
        df = pd.DataFrame(results)
        filepath = os.path.join(
            app.config["UPLOAD_FOLDER"], f"export_{session_id}.xlsx"
        )
        df.to_excel(filepath, index=False)
        return send_file(
            filepath, as_attachment=True, download_name="bloom_results.xlsx"
        )
    elif format_type == "word":
        doc = docx.Document()
        doc.add_heading("Classification Results", 0)
        for r in results:
            doc.add_heading(f"Q: {r['question']}", level=2)
            doc.add_paragraph(
                f"Bloom: {r['bloom_level']} | Diff: {r['difficulty']} | Conf: {
                    r['confidence']
                }%"
            )
            doc.add_paragraph(f"Exp: {r.get('explanation', '')}")
            doc.add_paragraph("-" * 20)
        filepath = os.path.join(
            app.config["UPLOAD_FOLDER"], f"export_{session_id}.docx"
        )
        doc.save(filepath)
        return send_file(
            filepath, as_attachment=True, download_name="bloom_results.docx"
        )
    elif format_type == "pdf":
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("helvetica", size=12)
        pdf.cell(
            0,
            10,
            text="Classification Results",
            new_x="LMARGIN",
            new_y="NEXT",
            align="C",
        )
        for r in results:
            q_safe = str(r["question"]).encode("latin-1", "replace").decode("latin-1")
            exp_safe = (
                str(r.get("explanation", ""))
                .encode("latin-1", "replace")
                .decode("latin-1")
            )
            text = f"Q: {q_safe}\nBloom: {r['bloom_level']} | Diff: {
                r['difficulty']
            } | Conf: {r['confidence']}%\nExp: {exp_safe}\n---"
            pdf.multi_cell(0, 10, text=text, new_x="LMARGIN", new_y="NEXT")
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], f"export_{session_id}.pdf")
        pdf.output(filepath)
        return send_file(
            filepath, as_attachment=True, download_name="bloom_results.pdf"
        )
    elif format_type == "ppt":
        prs = pptx.Presentation()

        for r in results:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Question"
            tf = content.text_frame
            tf.text = r["question"]

            tf.add_paragraph().text = f"Bloom Level: {r['bloom_level']}"
            tf.add_paragraph().text = f"Difficulty: {r['difficulty']}"
            tf.add_paragraph().text = f"Confidence: {r['confidence']}%"

        filepath = os.path.join(
            app.config["UPLOAD_FOLDER"], f"export_{session_id}.pptx"
        )
        prs.save(filepath)
        return send_file(
            filepath, as_attachment=True, download_name="bloom_results.pptx"
        )

    return jsonify({"error": "Invalid export format"}), 400


if __name__ == "__main__":
    DEBUG_MODE = False
    is_child = os.environ.get("WERKZEUG_RUN_MAIN") == "true"

    if DEBUG_MODE and not is_child:
        print("[PROCESS] Flask reloader parent process active. Skipping model loading.")
    else:
        print("[PROCESS] Flask serving process active. Spawning model loading thread...")
        threading.Thread(target=load_models, daemon=True).start()

    app.run(debug=DEBUG_MODE, port=5000)
